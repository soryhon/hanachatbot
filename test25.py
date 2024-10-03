import os
import openai
import streamlit as st
from langchain.prompts import PromptTemplate
from langchain import LLMChain
import requests
import PyPDF2
import pandas as pd
import docx
import pptx
from PIL import Image
from io import BytesIO
import base64  # base64 인코딩을 위해 추가
from langchain.chat_models import ChatOpenAI

# 전역변수로 프롬프트 저장
global_generated_prompt = ""

# 페이지 너비를 전체 화면으로 설정
st.set_page_config(layout="wide")

# GitHub API 요청을 처리하는 함수
def get_github_files(repo, branch, token, folder_name=None):
    # GitHub API 호출 (예시로 파일 리스트 반환)
    if folder_name:
        return [f"{folder_name}/file1.txt", f"{folder_name}/file2.txt"]
    else:
        return ["file1.txt", "file2.txt"]

# GitHub에 파일 업로드 함수 (폴더 없으면 생성)
def upload_file_to_github(repo, folder, file_name, content, token, branch="main", sha=None):
    url = f"https://api.github.com/repos/{repo}/contents/{folder}/{file_name}"
    headers = {"Authorization": f"token {token}"}
    
    # 파일 콘텐츠를 base64로 인코딩
    content_base64 = base64.b64encode(content).decode('utf-8')

    data = {
        "message": f"Upload {file_name} to {folder}",
        "content": content_base64,  # base64 인코딩된 콘텐츠
        "branch": branch
    }
    if sha:  # 덮어쓰기인 경우 SHA 값 전달
        data["sha"] = sha

    response = requests.put(url, headers=headers, json=data)
    if response.status_code in [200, 201]:
        st.success(f"'{file_name}' 파일이 성공적으로 업로드되었습니다!")
    else:
        st.error(f"파일 업로드 실패: {response.status_code}")

# GitHub에 폴더가 존재하는지 확인하는 함수
def get_file_sha(repo, path, token, branch="main"):
    url = f"https://api.github.com/repos/{repo}/contents/{path}?ref={branch}"
    headers = {"Authorization": f"token {token}"}
    response = requests.get(url, headers=headers)
    if response.status_code == 200:
        return response.json().get("sha", None)
    else:
        return None

# GitHub에서 파일을 다운로드하는 함수
def get_file_from_github(repo, branch, filepath, token):
    url = f"https://api.github.com/repos/{repo}/contents/{filepath}?ref={branch}"
    headers = {"Authorization": f"token {token}"}
    response = requests.get(url, headers=headers)
    
    if response.status_code == 200:
        return BytesIO(requests.get(response.json()['download_url']).content)
    else:
        st.error(f"{filepath} 파일을 가져오지 못했습니다.")
        return None

# 다양한 파일 형식에서 데이터를 추출하는 함수
def extract_data_from_file(file_content, file_type):
    if file_type == 'pdf':
        return extract_text_from_pdf(file_content)
    elif file_type == 'xlsx':
        return extract_text_from_excel(file_content)
    elif file_type == 'csv':
        return extract_text_from_csv(file_content)
    elif file_type == 'docx':
        return extract_text_from_word(file_content)
    elif file_type == 'pptx':
        return extract_text_from_ppt(file_content)
    elif file_type in ['png', 'jpg', 'jpeg']:
        return extract_text_from_image(file_content)
    else:
        st.error(f"{file_type} 형식은 지원되지 않습니다.")
        return None

# PDF 파일에서 텍스트 추출
def extract_text_from_pdf(file_content):
    reader = PyPDF2.PdfReader(file_content)
    text = ''
    for page in range(len(reader.pages)):
        text += reader.pages[page].extract_text()
    return text

# 엑셀 파일에서 텍스트 추출 (시트 정보를 정확하게 가져오도록 수정)
def extract_text_from_excel(file_content):
    try:
        excel_data = pd.read_excel(file_content, sheet_name=None)  # 모든 시트를 불러옴
        all_data = {}

        for sheet_name, data in excel_data.items():
            all_data[sheet_name] = data

        return all_data  # 각 시트 이름과 데이터가 포함된 딕셔너리 반환
    except Exception as e:
        st.error(f"엑셀 파일을 불러오는 중 오류가 발생했습니다: {str(e)}")
        return None

# CSV 파일에서 텍스트 추출
def extract_text_from_csv(file_content):
    csv_data = pd.read_csv(file_content)
    return csv_data

# 워드 파일에서 텍스트 추출
def extract_text_from_word(file_content):
    doc = docx.Document(file_content)
    return '\n'.join([para.text for para in doc.paragraphs])

# PPT 파일에서 텍스트 추출
def extract_text_from_ppt(file_content):
    presentation = pptx.Presentation(file_content)
    text = ''
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + '\n'
    return text

# 이미지에서 텍스트 추출 (OCR)
def extract_text_from_image(file_content):
    image = Image.open(file_content)
    return "이미지에서 텍스트를 추출하는 기능은 구현되지 않았습니다."

# LLM을 통해 프롬프트와 파일을 전달하고 응답을 받는 함수
def run_llm_with_file_and_prompt(api_key, title, request, file_data):
    global global_generated_prompt
    openai.api_key = api_key

    # file_data_str 변수: 파일 데이터를 텍스트 형태로 변환하여 LLM에 전달
    if isinstance(file_data, pd.DataFrame):  # 엑셀, CSV의 경우 DataFrame을 텍스트로 변환
        file_data_str = file_data.to_string()
    elif isinstance(file_data, dict):  # 여러 시트를 가져온 경우
        file_data_str = "\n".join([f"시트 이름: {sheet_name}\n{data.to_string()}" for sheet_name, data in file_data.items()])
    else:
        file_data_str = str(file_data)

    # 프롬프트 템플릿 구성
    generated_prompt = f"""
    보고서 제목은 '{title}'로 하고, 아래의 파일 데이터를 분석하여 '{request}'를 요구 사항을 만족할 수 있도록 최적화된 보고서를 완성해.
    표로 표현 할 때는 table 태그 형식으로 구현해야 한다. th과 td 태그는 border는 사이즈 1이고 색상은 검정색으로 구성한다.
    표의 첫번째 행은 타이틀이 이므로 th태그로 구현하고 가운데 정렬, bold처리 해야 한다. 바탕색은 '#E7E6E6' 이어야 한다.
    예시와 같은 구조로 구성한다. 보고서 제목은 앞에 순번을 표시하고 바로 아래 요구 사항에 맞는 내용을 이어서 보여줘야 한다.
    예시 : '\r\n 1. 보고서 제목\r\n보고서 내용'
    파일 데이터: {file_data_str}
    """

    # 전역변수에 프롬프트 저장
    global_generated_prompt = generated_prompt

    # PromptTemplate 설정
    prompt_template = PromptTemplate(
        template=generated_prompt,
        input_variables=[]
    )

    # LangChain의 LLMChain 사용
    llm = ChatOpenAI(model_name="gpt-4o")
    chain = LLMChain(llm=llm, prompt=prompt_template)
    
    # LLM에 프롬프트를 전달하고 응답 받기
    response = chain.run({})

    return response

# 1 프레임
# 1. GitHub 정보 저장 및 OpenAI API 키 저장
st.subheader("1. GitHub 정보 저장 및 OpenAI API 키 저장")

col1, col2 = st.columns(2)  # 두 개의 컬럼으로 나눔

# GitHub 정보 저장
with col1:
    st.subheader("GitHub 정보 입력")
    repo_input = st.text_input("GitHub 저장소 (owner/repo 형식)", value="")  # 공백으로 설정
    branch_input = st.text_input("GitHub 브랜치", value="main")
    token_input = st.text_input("GitHub Token", type="password")

    if st.button("GitHub 정보 저장"):
        if repo_input and branch_input and token_input:
            st.session_state["github_repo"] = repo_input
            st.session_state["github_branch"] = branch_input
            st.session_state["github_token"] = token_input
            st.success("GitHub 정보가 성공적으로 저장되었습니다!")
        else:
            st.error("모든 GitHub 정보를 입력해야 합니다!")

# OpenAI API 키 저장
with col2:
    st.subheader("OpenAI API 키 저장")
    api_key_input = st.text_input("OpenAI API 키를 입력하세요", type="password")
    if st.button("API 키 저장"):
        if api_key_input:
            st.session_state["api_key"] = api_key_input
            os.environ["OPENAI_API_KEY"] = api_key_input
            st.success("API 키가 성공적으로 저장되었습니다!")
        else:
            st.error("API 키를 입력해야 합니다!")

# 2 프레임
# 2. 파일 업로드

# GitHub 정보가 저장되지 않은 경우 업로드 금지
if not st.session_state.get("github_token") or not st.session_state.get("github_repo"):
    st.warning("GitHub 정보가 저장되기 전에는 파일 업로드를 할 수 없습니다. 먼저 GitHub 정보를 입력해 주세요.")
else:
    with st.expander("파일 업로드", expanded=True):
        uploaded_files = st.file_uploader("파일을 여러 개 드래그 앤 드롭하여 업로드하세요.", accept_multiple_files=True)

        if uploaded_files:
            for uploaded_file in uploaded_files:
                file_content = uploaded_file.read()
                file_name = uploaded_file.name
                folder_name = 'uploadFiles'

                sha = get_file_sha(st.session_state['github_repo'], f"{folder_name}/{file_name}", st.session_state['github_token'], branch=st.session_state['github_branch'])

                if sha:
                    st.warning(f"'{file_name}' 파일이 이미 존재합니다. 덮어쓰시겠습니까?")
                    col1, col2 = st.columns(2)

                    # 덮어쓰기 버튼
                    with col1:
                        if st.button(f"'{file_name}' 덮어쓰기"):
                            upload_file_to_github(st.session_state['github_repo'], folder_name, file_name, file_content, st.session_state['github_token'], branch=st.session_state['github_branch'], sha=sha)
                            st.success(f"'{file_name}' 파일이 성공적으로 덮어쓰기 되었습니다.")
                            uploaded_files = None  # 업로드 후 파일 리스트 초기화
                            break  # 업로드 완료 후 루프 종료

                    # 취소 버튼
                    with col2:
                        if st.button("취소"):
                            st.info("덮어쓰기가 취소되었습니다.")
                            uploaded_files = None  # 취소 후 파일 리스트 초기화
                            break  # 루프 종료
                else:
                    upload_file_to_github(st.session_state['github_repo'], folder_name, file_name, file_content, st.session_state['github_token'])
                    st.success(f"'{file_name}' 파일이 성공적으로 업로드되었습니다.")
                    uploaded_files = None  # 업로드 후 파일 리스트 초기화

# 3 프레임: 작성 보고서 요청사항 테이블
st.subheader("3. 작성 보고서 요청사항 및 실행 버튼")

# 3 프레임 안에 두 개의 영역을 나눔 (가로 80%, 20%)
col1, col2 = st.columns([0.8, 0.2])

with col1:
    # 요청사항 테이블 및 요청사항 관리 버튼 포함
    with st.expander("요청사항 리스트", expanded=True):
        if 'rows' not in st.session_state:
            # 페이지 접속 시 기본 행 1개 생성
            st.session_state['rows'] = [{"제목": "", "요청": "", "데이터": "", "checked": False}]

        rows = st.session_state['rows']  # 세션 상태에 저장된 행 목록을 사용
        checked_rows = []  # 체크된 행들을 저장하기 위한 리스트

        # 각 요청사항 테이블(순번마다 하나씩)
        for idx, row in enumerate(rows):
            with st.container():
                col1_1, col2_1 = st.columns([0.05, 0.95])
                with col1_1:
                    # 체크박스: 요청사항 + 순번
                    row_checked = st.checkbox(f"", key=f"row_checked_{idx}", value=row.get("checked", False), disabled=(idx == 0))  # 최초 요청사항1의 체크박스는 비활성화
                with col2_1:
                    st.markdown(f"#### 요청사항 {idx+1}")

                    # 제목
                    row['제목'] = st.text_input(f"제목 (요청사항 {idx+1})", row['제목'], key=f"title_{idx}")

                    # 요청
                    row['요청'] = st.text_input(f"요청 (요청사항 {idx+1})", row['요청'], key=f"request_{idx}")

                    # GitHub 파일 목록 불러오기
                    file_list = []
                    if st.session_state.get('github_repo') and st.session_state.get('github_token'):
                        file_list = get_github_files(st.session_state['github_repo'], st.session_state['github_token'], branch=st.session_state.get('github_branch', 'main'))
                    
                    selected_file = st.selectbox(f"파일 선택 (요청사항 {idx+1})", options=file_list, key=f"file_select_{idx}")

                    # 선택된 파일 서버 경로 저장
                    if st.button(f"선택 (요청사항 {idx+1})") and selected_file:
                        server_path = get_file_server_path(st.session_state['github_repo'], st.session_state.get('github_branch', 'main'), selected_file)
                        row['데이터'] = server_path
                        st.success(f"선택한 파일: {selected_file}\n서버 경로: {server_path}")

                    # 데이터 (선택된 파일의 경로)
                    st.text_input(f"데이터 (요청사항 {idx+1})", row['데이터'], disabled=True, key=f"data_{idx}")

                if row_checked:
                    checked_rows.append(idx)
                    row["checked"] = True
                else:
                    row["checked"] = False

        # 행 추가/삭제/새로고침 버튼
        col1_1, col1_2, col1_3 = st.columns([0.33, 0.33, 0.33])

        with col1_1:
            # 행 추가
            if st.button("행 추가"):
                st.session_state['rows'].append({"제목": "", "요청": "", "데이터": "", "checked": False})
                st.experimental_rerun()  # 새로고침을 통해 페이지 갱신

        with col1_2:
            # 행 삭제
            if st.button("행 삭제"):
                if checked_rows:
                    # 선택된 요청사항 삭제
                    st.session_state['rows'] = [row for idx, row in enumerate(rows) if idx not in checked_rows]
                    st.success(f"체크된 {len(checked_rows)}개의 요청사항이 삭제되었습니다.")
                    st.experimental_rerun()  # 새로고침을 통해 페이지 갱신
                else:
                    st.warning("삭제할 요청사항을 선택해주세요.")

        with col1_3:
            # 새로고침
            if st.button("새로고침"):
                st.experimental_rerun()  # 페이지를 새로고침

# 실행 버튼 영역
with col2:
    st.write(" ")
    st.write(" ")
    # 실행 버튼
    if st.button("실행"):
        if not st.session_state.get("api_key"):
            st.error("먼저 OpenAI API 키를 입력하고 저장하세요!")
        elif not st.session_state['rows'] or all(not row["제목"] or not row["요청"] for row in st.session_state['rows']):
            st.error("요청사항의 제목과 요청을 모두 입력해야 합니다!")
        else:
            # 요청사항을 기반으로 보고서를 생성
            response = run_llm_with_file_and_prompt(
                st.session_state["api_key"], 
                st.session_state['rows'][0]["제목"], 
                st.session_state['rows'][0]["요청"], 
                st.session_state['rows'][0]["데이터"]
            )
            st.session_state["response"] = response

# 4 프레임
# 4. 결과 보고서
st.subheader("4. 결과 보고서")

# 전달 프롬프트 출력
st.text_area("전달된 프롬프트:", value=global_generated_prompt, height=150)

if "response" in st.session_state:
    # 응답 텍스트 출력
    st.text_area("응답:", value=st.session_state["response"], height=300)
    
    # 응답 데이터를 HTML 형식으로 표시
    st.components.v1.html(st.session_state["response"], height=600, scrolling=True)
