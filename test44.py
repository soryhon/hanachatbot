import os
import openai
import streamlit as st
from langchain.prompts import PromptTemplate
from langchain import LLMChain
import requests
import PyPDF2
import pandas as pd
import docx
import pptx
from PIL import Image
from io import BytesIO
import base64
import urllib.parse
from openai.error import RateLimitError
from langchain.chat_models import ChatOpenAI
import time
import json
import openpyxl
from openpyxl.utils import get_column_letter
import re
import tempfile
import datetime

# Backend 기능 구현 시작 ---

# 전역변수로 프롬프트 및 파일 데이터 저장
global_generated_prompt = []

# GitHub 정보 및 OpenAI API 키 자동 설정 또는 입력창을 통해 설정
def load_env_info():
    json_data = '''
    {
        "github_repo": "soryhon/hanachatbot",
        "github_branch": "main"
    }
    '''
    
    # JSON 데이터에서 정보 추출
    github_info = json.loads(json_data)
    github_repo = github_info['github_repo']
    github_branch = github_info['github_branch']
    
    # 환경 변수에서 GitHub Token 및 OpenAI API Key 가져오기
    github_token = os.getenv("GITHUB_TOKEN")
    openai_api_key = os.getenv("OPENAI_API_KEY")

    github_set = False
    openai_set = False

    # 입력창을 가로로 배치 (각각 50%의 너비)
    col1, col2 = st.columns(2)

    with col1:
        if not github_token:
            github_token = st.text_input("GitHub Token을 입력하세요", type="password", key="github_token_input")
        else:
            github_set = True
            #st.success("GitHub 정보가 자동으로 설정되었습니다!")
        st.session_state["github_token"] = github_token

    with col2:
        if not openai_api_key:
            openai_api_key = st.text_input("OpenAI API 키를 입력하세요", type="password", key="openai_api_key_input")
        else:
            openai_set = True
            #st.success("OpenAI API 키가 자동으로 설정되었습니다!")
        st.session_state["openai_api_key"] = openai_api_key

    # GitHub 저장소 정보 세션에 저장
    st.session_state["github_repo"] = github_repo
    st.session_state["github_branch"] = github_branch

    # GitHub 정보가 설정되었는지 확인하고 세션 상태 반영
    return github_set

# GitHub에서 unloadFiles 하위의 폴더 리스트를 가져오는 함수
def get_folder_list_from_github(repo, branch, token, base_folder='uploadFiles'):
    url = f"https://api.github.com/repos/{repo}/contents/{base_folder}?ref={branch}"
    headers = {"Authorization": f"token {token}"}
    response = requests.get(url, headers=headers)
    if response.status_code == 200:
        folders = [item['name'] for item in response.json() if item['type'] == 'dir']
        return folders
    else:
        st.error("폴더 리스트를 가져오지 못했습니다. 저장소 정보나 토큰을 확인하세요.")
        return []

# GitHub에 새로운 폴더를 생성하는 함수
def create_new_folder_in_github(repo, folder_name, token, branch='main'):
    base_folder = "uploadFiles"
    folder_path = f"{base_folder}/{folder_name}/.gitkeep"  # Git에서 빈 폴더를 유지하는 방법 중 하나인 .gitkeep 파일 사용
    url = f"https://api.github.com/repos/{repo}/contents/{folder_path}"
    headers = {"Authorization": f"token {token}"}
    
    data = {
        "message": f"Create folder {folder_name}",
        "content": base64.b64encode(b'').decode('utf-8'),  # 빈 파일로 폴더 생성
        "branch": branch
    }
    
    response = requests.put(url, json=data, headers=headers)
    
    if response.status_code in [200, 201]:
        #st.success(f"'{folder_name}' 폴더가 성공적으로 생성되었습니다.")
        return True
    elif response.status_code == 422:
        st.warning("이미 존재하는 폴더입니다.")
        return False
    else:
        st.error(f"폴더 생성 실패: {response.status_code}")
        return False
      
# 다양한 파일 형식에서 데이터를 추출하는 함수
def extract_data_from_file(file_content, file_type):
    if file_content is None:
        st.error("파일 내용을 가져오지 못했습니다.")
        return None

    if file_type == 'pdf':
        return extract_text_from_pdf(file_content)
    elif file_type == 'csv':
        return extract_text_from_csv(file_content)
    elif file_type == 'docx':
        return extract_text_from_word(file_content)
    elif file_type == 'pptx':
        return extract_text_from_ppt(file_content)
    elif file_type in ['png', 'jpg', 'jpeg']:
        return extract_text_from_image(file_content)
    elif file_type == 'txt':
        return extract_text_from_txt(file_content)
    elif file_type == 'log':
        return extract_text_from_log(file_content)
    elif file_type == 'wav':  # 음성 파일 추가
        return extract_text_from_audio(file_content, file_type)    
    else:
        st.error(f"{file_type} 형식은 지원되지 않습니다.")
        return None

# PDF 파일에서 텍스트 추출
def extract_text_from_pdf(file_content):
    reader = PyPDF2.PdfReader(file_content)
    text = ''
    for page in range(len(reader.pages)):
        text += reader.pages[page].extract_text()
    return text

# CSV 파일에서 텍스트 추출
def extract_text_from_csv(file_content):
    csv_data = pd.read_csv(file_content)
    return csv_data

# 워드 파일에서 텍스트 추출
def extract_text_from_word(file_content):
    doc = docx.Document(file_content)
    return '\n'.join([para.text for para in doc.paragraphs])

# PPT 파일에서 텍스트 추출
def extract_text_from_ppt(file_content):
    presentation = pptx.Presentation(file_content)
    text = ''
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + '\n'
    return text

# 텍스트 파일에서 텍스트 추출
def extract_text_from_txt(file_content):
    try:
        # BytesIO 객체의 내용을 문자열로 변환
        if isinstance(file_content, BytesIO):
            return file_content.read().decode("utf-8")
        elif isinstance(file_content, str):
            return file_content
        else:
            st.error("알 수 없는 파일 형식입니다.")
            return None
    except Exception as e:
        st.error(f"txt 파일에서 텍스트를 추출하는 중 오류가 발생했습니다: {str(e)}")
        return None

# 로그 파일에서 텍스트 추출
def extract_text_from_log(file_content):
    try:
        # BytesIO 객체의 내용을 문자열로 변환
        if isinstance(file_content, BytesIO):
            return file_content.read().decode("utf-8")
        elif isinstance(file_content, str):
            return file_content
        else:
            st.error("알 수 없는 파일 형식입니다.")
            return None
    except Exception as e:
        st.error(f"log 파일에서 텍스트를 추출하는 중 오류가 발생했습니다: {str(e)}")
        return None

    
# 이미지에서 텍스트 추출 (OCR)
def extract_text_from_image(file_content):
    image = Image.open(file_content)
    return "이미지에서 텍스트를 추출하는 기능은 구현되지 않았습니다."


# GitHub에 폴더가 존재하는지 확인하고 없으면 생성하는 함수
def create_github_folder_if_not_exists(repo, folder_name, token, branch='main'):
    url = f"https://api.github.com/repos/{repo}/contents/{folder_name}?ref={branch}"
    headers = {"Authorization": f"token {token}"}
    response = requests.get(url, headers=headers)
    
    if response.status_code == 404:
        st.warning(f"'{folder_name}' 폴더가 존재하지 않아 생성 중입니다.")
        create_folder_url = f"https://api.github.com/repos/{repo}/contents/{folder_name}"
        data = {
            "message": f"Create {folder_name} folder",
            "content": base64.b64encode(b'').decode('utf-8'),  # 빈 파일로 폴더 생성
            "branch": branch
        }
        requests.put(create_folder_url, json=data, headers=headers)
        #st.success(f"'{folder_name}' 폴더가 성공적으로 생성되었습니다.")
    # 폴더가 이미 존재할 경우 메시지를 표시하지 않음

# GitHub API 요청을 처리하는 함수 (파일 목록을 가져옴)
def get_github_files(repo, branch, token):
    # 보고서명 리스트에서 선택한 폴더가 upload_folder에 저장됨
    folder_name = st.session_state.get('upload_folder', 'uploadFiles')
    
    # upload_folder 하위 폴더 내의 파일을 가져옴
    create_github_folder_if_not_exists(repo, folder_name, token, branch)
    url = f"https://api.github.com/repos/{repo}/git/trees/{branch}?recursive=1"
    headers = {"Authorization": f"token {token}"}
    response = requests.get(url, headers=headers)
    
    if response.status_code == 200:
        files = [item['path'] for item in response.json().get('tree', []) if item['type'] == 'blob' and item['path'].startswith(folder_name)]
        return files
    else:
        st.error("GitHub 파일 목록을 가져오지 못했습니다. 저장소 정보나 토큰을 확인하세요.")
        return []

# GitHub에서 파일의 SHA 값을 가져오는 함수
def get_file_sha(repo, file_path, token, branch='main'):
    encoded_file_path = urllib.parse.quote(file_path)
    url = f"https://api.github.com/repos/{repo}/contents/{encoded_file_path}?ref={branch}"
    headers = {"Authorization": f"token {token}"}
    response = requests.get(url, headers=headers)
    
    if response.status_code == 200:
        return response.json().get('sha', None)
    else:
        return None

# GitHub에 파일 업로드 함수 (덮어쓰기 포함)
def upload_file_to_github(repo, folder_name, file_name, file_content, token, branch='main', sha=None):
    create_github_folder_if_not_exists(repo, folder_name, token, branch)
    encoded_file_name = urllib.parse.quote(file_name)
    url = f"https://api.github.com/repos/{repo}/contents/{folder_name}/{encoded_file_name}"
    headers = {
        "Authorization": f"token {token}",
        "Content-Type": "application/json"
    }

    content_encoded = base64.b64encode(file_content).decode('utf-8')

    data = {
        "message": f"Upload {file_name}",
        "content": content_encoded,
        "branch": branch
    }

    if sha:
        data["sha"] = sha

    response = requests.put(url, json=data, headers=headers)

    if response.status_code in [200, 201]:
        st.success(f"{file_name} 파일이 성공적으로 업로드(덮어쓰기) 되었습니다.")
    else:
        st.error(f"파일 업로드에 실패했습니다: {response.status_code}")
        st.error(response.json())

# GitHub에서 파일을 다운로드하는 함수
def get_file_from_github(repo, branch, filepath, token):
    encoded_filepath = urllib.parse.quote(filepath)
    url = f"https://api.github.com/repos/{repo}/contents/{encoded_filepath}?ref={branch}"
    headers = {"Authorization": f"token {token}"}
    response = requests.get(url, headers=headers)
    
    if response.status_code == 200:
        return BytesIO(requests.get(response.json()['download_url']).content)
    else:
        st.error(f"{filepath} 파일을 가져오지 못했습니다. 상태 코드: {response.status_code}")
        return None

# 엑셀 시트에서 셀 스타일 및 병합 정보 추출
def extract_cell_style_and_merged(ws):
    style_dict = {}
    merged_cells = ws.merged_cells.ranges
    
    for row in ws.iter_rows():
        for cell in row:
            cell_style = {
                "alignment": cell.alignment.horizontal if cell.alignment else 'left',
                "font_bold": cell.font.bold if cell.font else False,
                "border": bool(cell.border and (cell.border.left or cell.border.right or cell.border.top or cell.border.bottom)),
                "value": "" if str(cell.value).startswith("Unnamed") else cell.value
            }
            style_dict[cell.coordinate] = cell_style
            
    return style_dict, merged_cells

# 엑셀 시트 데이터를 HTML로 변환하고 스타일 및 병합 적용
def convert_df_to_html_with_styles_and_merging(ws, df):
    # 1행의 'Unnamed: 숫자' 형식 값은 공백으로 처리
    df.columns = [re.sub(r'Unnamed: \d+', '', str(col)).strip() for col in df.columns]
    
    style_dict, merged_cells = extract_cell_style_and_merged(ws)
    df = df.fillna('')  # NaN 값을 공백으로 처리
    html = "<table class='table table-bordered' style='border-spacing: 0; border-collapse: collapse;'>\n"

    # 헤더 부분 (border 간격 없게 설정)
    html += "<thead>\n<tr>\n"
    for col in df.columns:
        html += f"<th style='text-align:center; font-weight:bold; background-color:#E7E6E6; border: 1px solid black;'>{col}</th>\n"
    html += "</tr>\n</thead>\n"

    # 병합 셀 정보 저장
    merged_cells_dict = {}
    for merged in merged_cells:
        min_row, min_col, max_row, max_col = merged.min_row, merged.min_col, merged.max_row, merged.max_col
        for row in range(min_row, max_row + 1):
            for col in range(min_col, max_col + 1):
                merged_cells_dict[(row, col)] = (min_row, min_col, max_row, max_col)

    # 데이터 부분
    html += "<tbody>\n"
    for row_idx, row in df.iterrows():
        html += "<tr>\n"
        for col_idx, value in enumerate(row):
            cell_ref = f"{get_column_letter(col_idx+1)}{row_idx+2}"  # 셀 참조 계산
            style = style_dict.get(cell_ref, {})
            alignment = style.get("alignment", "left")
            font_weight = "bold" if style.get("font_bold", False) else "normal"
            border = "1px solid black" if style.get("border", False) else "none"
            cell_coordinates = (row_idx + 2, col_idx + 1)

            if cell_coordinates in merged_cells_dict:
                min_row, min_col, max_row, max_col = merged_cells_dict[cell_coordinates]
                rowspan = max_row - min_row + 1
                colspan = max_col - min_col + 1
                if (row_idx + 2, col_idx + 1) == (min_row, min_col):
                    html += f"<td rowspan='{rowspan}' colspan='{colspan}' style='text-align:{alignment}; font-weight:{font_weight}; border:{border};'>{value}</td>\n"
            elif cell_coordinates not in merged_cells_dict:
                html += f"<td style='text-align:{alignment}; font-weight:{font_weight}; border:{border};'>{value}</td>\n"
        html += "</tr>\n"
    html += "</tbody>\n</table>"
    
    return html

# 엑셀 파일에서 시트 가져오기 및 데이터 추출
def extract_sheets_from_excel(file_content, selected_sheets):
    try:
        excel_data = pd.ExcelFile(file_content)
        all_sheets = excel_data.sheet_names
        
        # 사용자가 선택한 시트를 가져옴
        if selected_sheets == 'all':
            selected_sheets = all_sheets
        else:
            selected_sheets = [all_sheets[int(i)-1] for i in selected_sheets if int(i) <= len(all_sheets)]
        
        # 선택한 시트의 데이터를 하나로 합침
        data_dict = {}
        with pd.ExcelFile(file_content) as xls:
            for sheet in selected_sheets:
                df = pd.read_excel(xls, sheet_name=sheet)
                data_dict[sheet] = df
                
        return data_dict
    except Exception as e:
        st.error(f"엑셀 파일의 시트 데이터를 추출하는 중에 오류가 발생했습니다: {str(e)}")
        return None

# 숫자, ',' 및 '-'만 허용하는 함수
def validate_sheet_input(input_value):
    if all(c.isdigit() or c in ['-', ','] for c in input_value):
        return True
    return False

# 시트 선택 로직 추가
def handle_sheet_selection(file_content, sheet_count, idx):
    # 3개의 객체를 가로로 배치
    col1, col2 = st.columns([0.5, 0.5])
    with col1:
        st.text_input(f"시트 갯수_{idx}", value=f"{sheet_count}개", key=f"sheet_count_{idx}", disabled=True)  # 시트 갯수 표시 (비활성화)
    
    with col2:
        # st.session_state['rows']와 st.session_state['rows'][idx]['파일정보']가 유효한지 확인하여 값 설정
        sheet_selection_value = '1'
        if st.session_state.get('rows') and st.session_state['rows'][idx].get('파일') and st.session_state['rows'][idx].get('파일정보'):
            sheet_selection_value = st.session_state['rows'][idx]['파일정보']
            file_name = st.session_state['rows'][idx]['파일']
        
        sheet_selection = st.text_input(
            f"시트 선택_{idx}(예: 1-3, 5)", 
            value=sheet_selection_value, 
            key=f"sheet_selection_{idx}_{file_name}"
        )
    
        # 입력값 변경 시 세션에 저장
        if validate_sheet_input(sheet_selection):
            st.session_state['rows'][idx]['파일정보'] = sheet_selection
        else:
            st.error("잘못된 입력입니다. 숫자와 '-', ',' 만 입력할 수 있습니다.")
    return None

# 시트 선택 입력값을 분석하는 함수
def parse_sheet_selection(selection, sheet_count):
    selected_sheets = []

    try:
        if '-' in selection:
            start, end = map(int, selection.split('-'))
            if start <= end <= sheet_count:
                selected_sheets.extend(list(range(start, end+1)))
        elif ',' in selection:
            selected_sheets = [int(i) for i in selection.split(',') if 1 <= int(i) <= sheet_count]
        else:
            selected_sheets = [int(selection)] if 1 <= int(selection) <= sheet_count else []
    except ValueError:
        st.error("잘못된 시트 선택 입력입니다.")
        return None

    return selected_sheets

# 파일에서 데이터를 추출하고 요청사항 리스트에서 선택한 엑셀 파일의 시트를 보여주는 로직 수정
def handle_file_selection(file_path, file_content, file_type, idx):
    if file_type == 'xlsx':
        wb = openpyxl.load_workbook(file_content)
        sheet_count = len(wb.sheetnames)

        # 시트 선택 로직 처리
        handle_sheet_selection(file_content, sheet_count, idx)
        #return file_data_dict
    #else:
        #return extract_data_from_file(file_content, file_type)

# HTML 보고서 생성 함수 (배열에서 데이터 가져옴)
def generate_final_html_report(file_data):
    report_html = ""
    if "html_report" in st.session_state:
        report_html = st.session_state['html_report']
    if file_data:
        report_html += f"<div style='text-indent: 1px;'>\n{file_data}\n</div><p/.\n"
        st.session_state['html_report'] = report_html  # 최종 값을 세션 상태에 저장

# 엑셀 데이터 및 을 HTML로 변환하여 하나의 세트로 출력하는 함수
def generate_html_report_with_title(titles, data_dicts):
    report_html = ""
    
    for i, (title, data_dict) in enumerate(zip(titles, data_dicts), start=1):
        report_html += f"<h3>{i}. {title}</h3>\n"
        report_html += "<div style='text-indent: 20px;'>\n"
        
        for sheet_name, df in data_dict.items():
            wb = openpyxl.load_workbook(BytesIO(df))
            ws = wb[sheet_name]
            report_html += convert_df_to_html_with_styles_and_merging(ws, df)
        
        report_html += "</div>\n"
    
    return report_html

# LLM 연동 및 응답 처리 함수
def execute_llm_request(api_key, prompt):
    openai.api_key = api_key
    responses = []

    try:
        # 텍스트 길이 제한 확인 (예: 4000자로 제한)
        if len(prompt) > 4000:
            st.error("프롬프트 글자 수 초과로 LLM 연동에 실패했습니다.")
            return None

        # 프롬프트 템플릿 설정
        prompt_template = PromptTemplate(
            template=prompt,
            input_variables=[]
        )

        # LLM 모델 생성
        llm = ChatOpenAI(model_name="gpt-4o")
        chain = LLMChain(llm=llm, prompt=prompt_template)

        success = False
        retry_count = 0
        max_retries = 5  # 최대 재시도 횟수

        # 응답을 받을 때까지 재시도
        while not success and retry_count < max_retries:
            try:
                response = chain.run({})
                responses.append(response)
                success = True
            except RateLimitError:
                retry_count += 1
                st.warning(f"API 요청 한도를 초과했습니다. 10초 후 다시 시도합니다. 재시도 횟수: {retry_count}/{max_retries}")
                time.sleep(10)

            time.sleep(10)
    except Exception as e:
        st.error(f"LLM 실행 중 오류가 발생했습니다: {str(e)}")

    return responses

# LLM을 통해 프롬프트와 파일을 전달하고 응답을 받는 함수
def run_llm_with_file_and_prompt(api_key, titles, requests, file_data_str):
    global global_generated_prompt
    openai.api_key = api_key

    responses = []
    global_generated_prompt = []  # 프롬프트들을 담을 리스트 초기화

    try:
        # 요청사항 리스트 문자열 생성
        request_list_str = "\n".join([
            f"{i+1}.{title}의 항목 데이터에 대해 '{request}' 요청 사항을 만족하게 구성한다. 항목 데이터의 데이터 값은 중략하거나 누락되어서 안된다.\n"
            for i, (title, request) in enumerate(zip(titles, requests))
        ])

        # 프롬프트 텍스트 정의
        generated_prompt = f"""
        아애의 항목 데이터를 간결하고 깔끔한 보고서 작성을 위해 보고서 내용에 대해서 알기 쉽게 내용 요약하고 설명해야 한다.
        다음과 같은 조건에 모두 만족해야 한다.
        가. 아래의 항목 데이터를 분석하여 각 항목마다의 '요청사항' 리스트와 조건사항에 대해 모두 만족할 수 있도록 최적화된 보고서를 완성해.
        나. 항목 데이터 내 가장 첫번째 행은 각 보고서 항목에 타이틀이므로 순번과 문구를 그대로 유지해야 한다. 이 항목의 타이틀을 기준으로 각 항목 데이터를 분류하고 그에 맞는 요청사항을 반영해야 한다.
        다. 항목데이터 중 table 태그가 포함된 데이터는 엑셀에서 추출한 데이터로 표 형식으로 그대로 유지해야 한다. 
              table태그 데이터의 데이터 값은 중략하거나 누락되어서 안된다.보고서에 꼭 필요하니 최대한 그대로 구조와 데이터는 출력되게 하고 내용만 업데이트한다. 
        라. 표 형식의로 table태그로 답변 할 때는 th과 td 태그는 border는 사이즈 1이고 색상은 검정색으로 구성한다. table 태그 가로길이는 전체를 차지해야 한다.
        마. 표 형식으로 답변할 때는 반드시 모든 항목 데이터의 수정한 데이터 내용과 HTML 태그를 보완한 모든 데이터를 업그레이드한 데이터로 보여줘야 한다.
        바. 이외 table 태그가 포함 안된 데이터는 파일에서 텍스트를 추출한 데이터로 내용으로 보고서 양식에 맞춰 간결하고 깔끔하게 요약하고 html 형식으로 변환해야 한다.
        사. 답변할 때는 반드시 모든 항목 데이터의 수정한 데이터 내용과 HTML 형삭에 맞춰 답변한다. 문단마다 줄바꿈을 적용하여 br태그 활용하고 가시성 높게 특수기호를 활용하여 보고서 양식에 준하게 요약한 내용을 설명해줘야 한다.
        아. 항목 데이터를 업데이트 하는 결과 가장 먼저 나와야하고 위에는 그 어떤한 설명 내용이 응답하면 안 된다. 
        자. 그 다음 일정 간격을 두고 h3 태그를 활용해서 '✨AI 요약과 설명' 타이틀 추가하고 색상을 달리 구성한다. 너의 답변이라는 것을 표현하는 특수문자로 강조해.
               전달받은 보고서 전반적인 내용에 대해 너가 선정한 가장 좋은 방법으로 요약과 설명하고 그 내용을 HTML 형식으로 변환하여 답변해야 한다.
        차. '````', '````HTML' 이 문구들이 답변에 포함되지 않아야 한다.
        -요청사항
        [
            {request_list_str}
        ]
        -항목 데이터
        [
            {file_data_str}
        ]
        """
        
        # 텍스트 길이 제한 확인 (예: 1000000자로 제한)
        if len(generated_prompt) > 1000000:
            st.error("프롬프트 글자 수 초과로 LLM 연동에 실패했습니다.")
        else:
            global_generated_prompt.append(generated_prompt)
            prompt_template = PromptTemplate(
                template=generated_prompt,
                input_variables=[]
            )

            # LLM 모델 생성
            llm = ChatOpenAI(model_name="gpt-4o")
            chain = LLMChain(llm=llm, prompt=prompt_template)

            success = False
            retry_count = 0
            max_retries = 5  # 최대 재시도 횟수

            # 응답을 받을 때까지 재시도
            while not success and retry_count < max_retries:
                try:
                    response = chain.run({})
                    responses.append(response)
                    success = True
                except RateLimitError:
                    retry_count += 1
                    st.warning(f"API 요청 한도를 초과했습니다. 10초 후 다시 시도합니다. 재시도 횟수: {retry_count}/{max_retries}")
                    time.sleep(10)

                time.sleep(10)
    except Exception as e:
        st.error(f"LLM 실행 중 오류가 발생했습니다: {str(e)}")

    return responses


#st.session_state를 새로고침하는 함수
def refresh_page():
    if 'is_updating' not in st.session_state:
        st.session_state['is_updating'] = False
    elif st.session_state['is_updating']:
        st.session_state['is_updating'] = False
    else:
        st.session_state['is_updating'] = True

def init_session_state(check_value):
    folderlist_init_value = "보고서명을 선택하세요."
    templatelist_init_value = "불러올 보고서 양식을 선택하세요."
    if(check_value):
            st.session_state['rows'] = [
                {"제목": "", "요청": "", "파일": "", "데이터": "","파일정보":"1" }
                for _ in range(st.session_state['num_requests'])
            ]    
            st.session_state['html_report'] = ""
    else:
        if 'selected_folder_name' not in st.session_state:
            st.session_state['selected_folder_name'] = folderlist_init_value
        if 'folder_list_option' not in st.session_state:       
            st.session_state['folder_list_option'] = folderlist_init_value
        if 'selected_template_name' not in st.session_state:
            st.session_state['selected_template_name'] = templatelist_init_value
        if 'template_list_option' not in st.session_state:       
            st.session_state['template_list_option'] = templatelist_init_value
        if 'upload_folder' not in st.session_state:        
            st.session_state['upload_folder'] = "uploadFiles" 
        if 'selected_folder_index' not in st.session_state:    
            st.session_state['selected_folder_index'] = 0
        if 'selected_template_index' not in st.session_state:
            st.session_state['selected_template_index'] = 0
        if 'new_folder_text' not in st.session_state:    
            st.session_state['new_folder_text'] = ""
        if 'check_report' not in st.session_state:    
            st.session_state['check_report'] = True
        if 'check_upload' not in st.session_state:    
            st.session_state['check_upload'] = False        
        if 'check_request' not in st.session_state:    
            st.session_state['check_request'] = False
        if 'check_result' not in st.session_state:    
            st.session_state['check_result'] = False
        if 'check_count' not in st.session_state:    
            st.session_state['check_count'] = False
        if 'report_date_str' not in st.session_state: 
            st.session_state['report_date_str'] = ""
# HTML 파일을 저장하고 파일 경로를 반환하는 함수 (날짜 포함)
def save_html_response(html_content, folder_name, report_date_str):
    # 현재 시간을 'YYYYMMDDHHMMSS' 형식으로 가져오기
    #current_time = datetime.datetime.now().strftime("%Y%m%d%H%M%S")
    # HTML 파일명을 보고서명과 날짜로 설정
    file_name = f"{folder_name}_result_{report_date_str}.html"
    
    # HTML 파일 임시 경로에 저장
    temp_file_path = f"/tmp/{file_name}"
    with open(temp_file_path, 'w', encoding='utf-8') as f:
        f.write(html_content)
    
    return file_name, temp_file_path        

# GitHub에 폴더가 존재하는지 확인하는 함수
def check_and_create_github_folder(folder_name, repo, branch, token):
    url = f"https://api.github.com/repos/{repo}/contents/{folder_name}"
    headers = {"Authorization": f"token {token}"}

    # 폴더 존재 여부 확인
    response = requests.get(url, headers=headers)
    if response.status_code == 404:  # 폴더가 없으면 생성
        # 폴더 생성하기 위한 커밋 메시지와 파일 내용 설정
        data = {
            "message": f"Create {folder_name} folder",
            "content": "",  # GitHub에서는 폴더 자체를 직접 생성할 수 없으므로 빈 파일 생성으로 대체
            "branch": branch
        }
        # 빈 파일 생성 (ex: .gitkeep)
        file_url = f"{url}/.gitkeep"
        response = requests.put(file_url, headers=headers, data=json.dumps(data))
        if response.status_code == 201:
            st.success(f"{folder_name} 폴더가 생성되었습니다.")
        else:
            st.error(f"{folder_name} 폴더 생성 실패: {response.json()}")
    elif response.status_code == 200:
        #st.info(f"{folder_name} 폴더가 이미 존재합니다.")
        return None
    else:
        st.error(f"폴더 확인 중 오류 발생: {response.json()}")
        
# JSON 파일 저장 함수
def save_template_to_json():
    repo = st.session_state["github_repo"]
    branch = st.session_state["github_branch"]
    token = st.session_state["github_token"]

    # GitHub 토큰과 레포지토리 설정 확인
    if not token or not repo:
        st.error("GitHub 토큰이나 저장소 정보가 설정되지 않았습니다.")
        return
        
    # JSON 데이터 구조 생성
    template_data = {
        "selected_folder_name": st.session_state['selected_folder_name'],
        "num_requests": st.session_state['num_requests'],
        "rows": st.session_state['rows'],
        "rows_length": len(st.session_state['rows']),
        "timestamp": datetime.datetime.now().strftime("%Y%m%d%H%M%S")
    }

    # 파일명 생성
    folder_name = st.session_state['selected_folder_name']
    timestamp = template_data["timestamp"]
    json_file_name = f"{folder_name}_Template_{timestamp}.json"

    # GitHub 저장소 내 templateFiles 폴더 생성 및 파일 저장
    template_folder = "templateFiles"
    check_and_create_github_folder(template_folder, repo, branch, token)
   
    # 저장할 파일 경로
    json_file_path = f"{template_folder}/{json_file_name}"

    # JSON 파일을 Base64로 인코딩
    json_content = json.dumps(template_data, ensure_ascii=False, indent=4)
    json_base64 = base64.b64encode(json_content.encode('utf-8')).decode('utf-8')

    # GitHub에 파일 업로드
    url = f"https://api.github.com/repos/{repo}/contents/{json_file_path}"
    headers = {"Authorization": f"token {token}"}
    data = {
        "message": f"Add {json_file_name}",
        "content": json_base64,
        "branch": branch
    }

    response = requests.put(url, headers=headers, json=data)
    if response.status_code == 201:
        st.success(f"{json_file_name} 파일이 {template_folder} 폴더에 저장되었습니다.")
    else:
        st.error(f"파일 저장 실패: {response.json()}")

# GitHub에서 templateFiles 폴더 내의 JSON 파일 리스트를 가져오는 함수
def get_template_files_list(repo, branch, token):
    template_folder = "templateFiles"
    url = f"https://api.github.com/repos/{repo}/contents/{template_folder}?ref={branch}"
    headers = {"Authorization": f"token {token}"}
    response = requests.get(url, headers=headers)

    if response.status_code == 200:
        # JSON 파일만 필터링하여 리스트로 반환
        return [item['name'] for item in response.json() if item['name'].endswith('.json')]
    else:
        st.error("templateFiles 폴더의 파일 목록을 가져오지 못했습니다.")
        return []

# JSON 파일의 내용을 불러오는 함수
def load_template_from_github(repo, branch, token, file_name):
    template_folder = "templateFiles"
    json_file_path = f"{template_folder}/{file_name}"
    url = f"https://api.github.com/repos/{repo}/contents/{json_file_path}?ref={branch}"
    headers = {"Authorization": f"token {token}"}
    response = requests.get(url, headers=headers)

    if response.status_code == 200:
        file_content = base64.b64decode(response.json()['content'])
        return json.loads(file_content)
    else:
        st.error(f"{file_name} 파일을 가져오지 못했습니다.")
        return None

def apply_template_to_session_state(file_name):
    try:
        # 템플릿 JSON 파일 로드
        with open(file_name, 'r', encoding='utf-8') as f:
            template_data = json.load(f)
        
        # JSON 데이터에서 세션 상태 적용
        selected_folder_name = template_data.get('selected_folder_name', '')
        #num_requests = template_data.get('num_requests', 1)
        rows = template_data.get('rows', [])
        
        # 세션 상태에 값 저장

        st.session_state['selected_folder_name'] = selected_folder_name
        st.session_state['rows'] = rows
        st.session_state['is_updating'] = False
        st.session_state['upload_folder'] = f"uploadFiles/{selected_folder_name}"
        st.session_state['check_report'] = False
        st.session_state['check_upload'] = False
        st.session_state['check_request'] = True
        st.session_state['check_result'] = False
        st.session_state['selected_folder_index'] = 0
    
        # 'num_requests'는 직접 변경할 수 없으므로 Streamlit에서 제공하는 방법으로 값을 설정
        #if "num_requests" in st.session_state:
            #st.session_state["num_requests"] = num_requests

        
        # folder_list에서 selected_folder_name의 인덱스 찾기
        folder_list = st.session_state.get('folder_list_option', [])
        if selected_folder_name in folder_list:
            selected_index = folder_list.index(selected_folder_name)
            st.session_state['selected_folder_index'] = selected_index + 1
        
        # 엑셀 파일 처리: 파일 정보에 따라 시트 선택 입력창 추가
        for idx, row in enumerate(rows):
            selected_file_name = row.get("파일")
            file_info = row.get("파일정보", "1")
            
            if selected_file_name and selected_file_name.endswith('.xlsx'):
                # 시트 선택 로직 적용
                #file_content = get_file_from_github(
                    #st.session_state["github_repo"],
                    #st.session_state["github_branch"],
                    #selected_file_name,
                    #st.session_state["github_token"]
                #)
                #if file_content:
                    #handle_sheet_selection(file_content, len(openpyxl.load_workbook(file_content).sheetnames), idx)
                st.session_state['rows'][idx]['파일정보'] = file_info        
        st.success(f"'{selected_folder_name}' 양식을 불러오기 성공하였습니다.")
    
    except FileNotFoundError:
        st.error(f"파일 '{file_name}'을 찾을 수 없습니다.")
    except json.JSONDecodeError:
        st.error(f"'{file_name}' 파일을 파싱하는 중 오류가 발생했습니다. JSON 형식을 확인해주세요.")
    except Exception as e:
        st.error(f"템플릿 불러오기 중 오류가 발생했습니다: {str(e)}")
           
# Backend 기능 구현 끝 ---

# Frontend 기능 구현 시작 ---

# GitHub 정보가 있는지 확인하고 파일 업로드 객체를 출력
github_info_loaded = load_env_info()

# 업로드 가능한 파일 크기 제한 (100MB)
MAX_FILE_SIZE_MB = 100
MAX_FILE_SIZE_BYTES = 100 * 1024 * 1024

#Session_state 변수 초기화
folderlist_init_value = "보고서명을 선택하세요."
templatelist_init_value = "불러올 보고서 양식을 선택하세요."
# 세션 상태에 각 변수 없다면 초기화
init_session_state(False)
refresh_page()
     
    
# 1 프레임
# 보고서 타이틀
col1, col2 = st.columns([0.7,0.3])
with col1:
    st.markdown(
        "<p style='font-size:25px; font-weight:bold; color:#000000;'> 업무 보고서 자동 완성</p>",
        unsafe_allow_html=True
    )
with col2:
    st.markdown(
        "<p style='font-size:13px; font-weight:normal; color:#aaaaaa; margin-top:10px;'>by <b style='font-size:16px;color:#0099FF'>CheokCeock</b><b style='font-size:22px;color:#009999'>1</b> <b style='font-size:14px;'>prototype v.01</b></p>",
        unsafe_allow_html=True
    )

# 2 프레임
# 보고서명 및 폴더 선택, 새 폴더 만들기
if github_info_loaded:
    with st.expander("📝 보고서 선택", expanded=st.session_state['check_report']):
        tab1, tab2, tab3 = st.tabs(["등록된 보고서명 선택하기", "저장된 보고서 양식 불러오기","새로운 보고서명 만들기"])
        with tab1:
            col1, col2 = st.columns([0.21, 0.79])
            with col1:
                st.write("")
                st.markdown(
                    "<p style='font-size:14px; font-weight:bold; color:#000000;text-align:center;'>등록된<br/>보고서명 선택 </p>",
                    unsafe_allow_html=True
                )
            with col2:
                folder_list = get_folder_list_from_github(st.session_state['github_repo'], st.session_state['github_branch'], st.session_state['github_token'])
                # st.selectbox 위젯 생성 (이제 session_state['selected_folder'] 사용 가능)
    
                # 'selected_folder'가 folder_list에 있을 때만 index 설정
                selected_index = st.session_state['selected_folder_index']
                if st.session_state['selected_folder_name'] in folder_list:
                    selected_index = folder_list.index(st.session_state['selected_folder_name']) + 1
                #else:
                    #selected_index = 0  # 기본값으로 '주제를 선택하세요.' 선택
                st.session_state['selected_folder_index'] = selected_index
                st.session_state['folder_list_option'] = [] + folder_list
                # 폴더 선택 selectbox 생성 (새 폴더 추가 후, 선택값으로 설정)
                selected_folder = st.selectbox(
                    "등록된 보고서명 리스트",
                    options=st.session_state['folder_list_option'],  # 옵션 리스트에 새 폴더 반영
                    index=st.session_state['selected_folder_index'],  # 새로 선택된 폴더를 기본값으로 선택
                    key="selected_folder"
                )
                # 파일 업로드와 요청사항 리스트의 기본 폴더 설정
                if selected_folder != folderlist_init_value:
                    st.session_state['upload_folder'] = f"uploadFiles/{selected_folder}"
                    st.session_state['selected_folder_name'] = f"{selected_folder}"                  
                    st.session_state['check_report']=False
                    st.session_state['check_count']=True
                    st.session_state['selected_template_index'] = 0
                    refresh_page()
                    #st.success(f"[{selected_folder}] 보고서명이이 선택되었습니다.")
                #else:   
                    #st.warning("보고서명을 선택하세요.")
        with tab2:
            col1, col2 = st.columns([0.21, 0.79])
            with col1:
                st.write("")
                st.markdown(
                    "<p style='font-size:14px; font-weight:bold; color:#000000;text-align:center;'>저장된 보고서<br/>양식 불러오기</p>",
                    unsafe_allow_html=True
                )
            with col2:    
                repo = st.session_state["github_repo"]
                branch = st.session_state["github_branch"]
                token = st.session_state["github_token"]
                 # templateFiles 폴더 내 JSON 파일 리스트 가져오기
                template_files = get_template_files_list(repo, branch, token)
                
                if template_files:
                    # 'selected_template'가 template_files에 있을 때만 index 설정
                    #selected_temp_index = st.session_state['selected_template_index']
                    if st.session_state['selected_template_name'] in template_files:
                        selected_temp_index = template_files.index(st.session_state['selected_template_name']) + 1                        
                    else:
                        selected_temp_index = 0
                    st.session_state['selected_template_index'] = selected_temp_index    
                    #보고서 양식 파일 리스트
                    selected_template = st.selectbox(
                        "불러올 보고서 양식 파일 리스트", 
                        options=[templatelist_init_value] + template_files, 
                        index=st.session_state['selected_template_index'],
                        key="selected_template"
                    )
                    if selected_template != templatelist_init_value:
                        # 선택한 템플릿 불러오기
                        st.session_state['selected_template_name'] = selected_template
                        template_data = load_template_from_github(repo, branch, token, selected_template)
                        if template_data:
                            apply_template_to_session_state(f"templateFiles/{selected_template}")
                            #st.success(f"{selected_template} 양식을 성공적으로 불러왔습니다.")

        with tab3:
            col1, col2, col3 = st.columns([0.21, 0.5,0.29])
            with col1:
                st.write("")
                st.markdown(
                    "<p style='font-size:14px; font-weight:bold; color:#000000;text-align:center;'>새로운 보고서명<br/>만들기</p>",
                    unsafe_allow_html=True
                )
            with col2:
                new_folder_name = st.text_input("새로 등록할 보고서명 입력", max_chars=20, key="new_folder_name", value=st.session_state['new_folder_text'])
            with col3:
                st.markdown(
                    "<p style='font-size:18px; margin-top:27px;'></p>",
                    unsafe_allow_html=True
                )
                if st.button("보고서명 등록", key="new_folder", use_container_width=True):
                    if not new_folder_name:
                        st.warning("새로 등록할 보고서명을 입력하세요.")
                    elif new_folder_name in folder_list:
                        st.warning("이미 존재합니다.")
                    else:
                        # 폴더 생성 후 목록에 추가
                        folder_created = create_new_folder_in_github(st.session_state['github_repo'], new_folder_name, st.session_state['github_token'], st.session_state['github_branch'])
                        if folder_created:
                            folder_list.append(new_folder_name)  # 새 폴더를 리스트에 추가
                            #st.session_state['selected_folder_index'] = len(folder_list) - 1
                            #st.session_state['selected_template_index'] = 0
                            st.session_state['folder_list_option'] = [] + folder_list
                            st.session_state['upload_folder'] = f"uploadFiles/{new_folder_name}"
                            st.session_state['selected_folder_name'] = f"{new_folder_name}"
                            st.session_state['selected_template_name'] = templatelist_init_value
                            st.session_state['check_report']=False
                            st.session_state['check_count']=True
                            refresh_page()
                            init_session_state(True)
                            st.success("새로운 보고서명 등록 성공하였습니다.")            
        #st.markdown(
            #"<hr style='border-top:1px solid #dddddd;border-bottom:0px solid #dddddd;width:100%;padding:0px;margin:0px'></hr>",
            #unsafe_allow_html=True
        #)
      
else:
    st.warning("GitHub 정보가 설정되지 않았습니다. 먼저 GitHub Token을 입력해 주세요.")



col1, col2, col3 = st.columns([0.2, 0.6, 0.2])
with col1:
    st.write("")
with col2:   
    report_title = "작성할 보고서를 선택하세요."
    title_style="font-size:15px; font-weight:normal; color:#cccccc;border: 1px solid #dddddd;letter-spacing: 1px;"
    if 'selected_folder_name' in st.session_state:
        if st.session_state['selected_folder_name'] != folderlist_init_value:
            report_title = " [" + st.session_state['selected_folder_name'] + "] 보고서"
            title_style="font-size:20px; font-weight:bold; color:#000000;border: 0px solid #dddddd;letter-spacing: 4px;"
    st.markdown(
        f"<div style='text-align:center;{title_style};border-radius: 10px;width:100%;padding: 10px;margin-top:10px;margin-bottom:10px;'>{report_title}</div>",
        unsafe_allow_html=True
    )
   
with col3:
    st.write("")

# 3 프레임
#st.subheader("")
st.markdown(
    "<p style='font-size:18px; font-weight:bold; color:#007BFF;'>작성 보고서 요청사항</p>",
    unsafe_allow_html=True
)

# 4 프레임
# 파일 업로드
# 지원되는 파일 형식 리스트
supported_file_types = ['xlsx', 'pptx', 'docx', 'csv', 'png', 'jpg', 'jpeg', 'pdf', 'txt', 'log']

if github_info_loaded:
    with st.expander("⬆️ 데이터 파일 업로드", expanded=st.session_state['check_upload']):
        uploaded_files = st.file_uploader("파일을 여러 개 드래그 앤 드롭하여 업로드하세요. (최대 100MB)", accept_multiple_files=True)

        if uploaded_files:
            for uploaded_file in uploaded_files:
                file_type = uploaded_file.name.split('.')[-1].lower()

                if file_type not in supported_file_types:
                    st.error(f"지원하지 않는 파일입니다: {uploaded_file.name}")
                    continue

                if uploaded_file.size > MAX_FILE_SIZE_BYTES:
                    st.warning(f"'{uploaded_file.name}' 파일은 {MAX_FILE_SIZE_MB}MB 제한을 초과했습니다. 파일 크기를 줄이거나 GitHub에 직접 푸시하세요.")
                else:
                    file_content = uploaded_file.read()
                    file_name = uploaded_file.name
                    #folder_name = 'uploadFiles'
                    folder_name = st.session_state.get('upload_folder', 'uploadFiles')

                    sha = get_file_sha(st.session_state['github_repo'], f"{folder_name}/{file_name}", st.session_state['github_token'], branch=st.session_state['github_branch'])

                    if sha:
                        st.warning(f"'{file_name}' 파일이 이미 존재합니다. 덮어쓰시겠습니까?")
                        col1, col2 = st.columns(2)

                        with col1:
                            if st.button(f"'{file_name}' 덮어쓰기", key=f"overwrite_{file_name}"):
                                upload_file_to_github(st.session_state['github_repo'], folder_name, file_name, file_content, st.session_state['github_token'], branch=st.session_state['github_branch'], sha=sha)
                                st.success(f"'{file_name}' 파일이 성공적으로 덮어쓰기 되었습니다.")
                                uploaded_files = None
                                break

                        with col2:
                            if st.button("취소", key=f"cancel_{file_name}"):
                                st.info("덮어쓰기가 취소되었습니다.")
                                uploaded_files = None
                                break
                    else:
                        upload_file_to_github(st.session_state['github_repo'], folder_name, file_name, file_content, st.session_state['github_token'])
                        st.success(f"'{file_name}' 파일이 성공적으로 업로드되었습니다.")
                        uploaded_files = None
else:
    st.warning("GitHub 정보가 저장되기 전에는 파일 업로드를 할 수 없습니다. 먼저 GitHub 정보를 입력해 주세요.")

# 5 프레임
# 요청사항 갯수 설정 입력 및 버튼
with st.expander("⚙️ 요청사항 설정", expanded=st.session_state['check_count']):
    col1, col2, col3 = st.columns([0.21, 0.4, 0.39])
    with col1:
        st.markdown(
            "<p style='font-size:14px; font-weight:bold; color:#000000; margin-top:20px;text-align:center;'>요청사항 리스트<br/>갯수 설정</p>",
            unsafe_allow_html=True
        )
        
    with col2:
        # 요청사항 갯수 입력 (1-9)
        num_requests = st.number_input(
            "요청사항 갯수 입력창",
            min_value=1,
            max_value=9,
            value=1,
            step=1,
            key="num_requests"
        )
    
    with col3:
        st.markdown(
            "<p style='font-size:18px; margin-top:27px;'></p>",
            unsafe_allow_html=True
        )
        if st.button("설정", key="set_requests", use_container_width=True):
            # 설정 버튼 클릭 시 요청사항 리스트 초기화 및 새로운 요청사항 갯수 설정
            st.session_state['rows'] = [
                {"제목": "", "요청": "", "파일": "", "데이터": "", "파일정보": "1"}
                for _ in range(st.session_state['num_requests'])
            ]
            st.success(f"{st.session_state['num_requests']}개의 요청사항이 설정되었습니다.")
            st.session_state['check_request']=True
            st.session_state['check_count']=False
            refresh_page()
            init_session_state(True)
    

# 6 프레임
# 요청사항 리스트
with st.expander("✍️ 요청사항 리스트", expanded=st.session_state['check_request']):
    if 'rows' not in st.session_state:
        st.session_state['rows'] = [{"제목": "", "요청": "", "파일": "", "데이터": "", "파일정보":"1"}]

    rows = st.session_state['rows']
    checked_rows = []

    for idx, row in enumerate(rows):
        with st.container():
            #col1, col2 = st.columns([0.01, 0.99]) 
            #with col1:
                #row_checked = st.checkbox("", key=f"row_checked_{idx}", value=row.get("checked", False))  # 체크박스만 추가
                #st.write("")
            #with col2:
            st.markdown(
                f"<p style='font-size:16px; font-weight:bold; color:#000000; margin-top:5px;'>{idx+1}. 요청사항</p>",
                unsafe_allow_html=True
            )
        
            row['제목'] = st.text_input(f"제목 : '{idx+1}.요청사항'의 제목을 입력해주세요.", row['제목'], key=f"title_{idx}")
            row['요청'] = st.text_area(f"요청 : '{idx+1}.요청사항'의 요청할 내용을 입력해주세요.", row['요청'], key=f"request_{idx}")
     
            file_list = ['파일을 선택하세요.']
            if st.session_state.get('github_token') and st.session_state.get('github_repo'):
                file_list += get_github_files(st.session_state['github_repo'], st.session_state['github_branch'], st.session_state['github_token'])

            selected_file = st.selectbox(f"파일 선택 : '{idx+1}.요청사항'의 파일을 선택해주세요.", options=file_list, key=f"file_select_{idx}")

            if selected_file != '파일을 선택하세요.':
                st.session_state['rows'][idx]['파일'] = selected_file
           
                file_path = selected_file
                file_content = get_file_from_github(
                    st.session_state["github_repo"], 
                    st.session_state["github_branch"], 
                    file_path, 
                    st.session_state["github_token"]
                )

                if file_content:
                    file_type = file_path.split('.')[-1].lower()
                    
                    # 파일 형식 검증 (지원되는 파일만 처리)
                    if file_type not in supported_file_types:
                        st.error(f"지원하지 않는 파일입니다: {file_path}")
                        row['데이터'] = ""
                    else:      
                        handle_file_selection(file_path, file_content, file_type, idx)
                        
                else:
                    st.error(f"{selected_file} 파일을 GitHub에서 불러오지 못했습니다.")  
            st.text_input(f"{idx+1}.요청사항 선택한 파일", row['파일'], disabled=True, key=f"file_{idx}")
        
# 7 프레임
col1, col2, col3 = st.columns([0.2, 0.6, 0.2])

with col1:
    st.write("")
with col2:   

# 보고서 실행 버튼 클릭 시 함수 호출 수정
    if st.button("🚀 보고서 작성 실행", key="generate_report", use_container_width=True):
        st.session_state['check_result']=True
        st.session_state['check_report'] = False
        st.session_state['check_upload'] = False
        st.session_state['check_count'] = False
        if not st.session_state.get("openai_api_key"):
            st.error("먼저 OpenAI API 키를 입력하고 저장하세요!")
        elif not st.session_state['rows'] or all(not row["제목"] or not row["요청"] or not row["파일"] for row in st.session_state['rows']):
            st.error("요청사항의 제목, 요청, 파일을 모두 입력해야 합니다!")
        else:
            with st.spinner('요청사항과 파일 데이터을 추출 중입니다...'):
        
                # 파일 데이터 가져와서 HTML 보고서 생성
                #file_data_list = []
                html_viewer_data = ""
                for idx, row in enumerate(st.session_state['rows']):
                    file_path = st.session_state['rows'][idx]['파일']
                    file_content = get_file_from_github(st.session_state["github_repo"], st.session_state["github_branch"], file_path, st.session_state["github_token"])
                    file_type = file_path.split('.')[-1].lower()
                    report_html = ""
                    if file_content:
                        if file_type == 'xlsx':
                            selected_sheets = parse_sheet_selection(row['파일정보'], len(openpyxl.load_workbook(file_content).sheetnames))
                            file_data_dict = extract_sheets_from_excel(file_content, selected_sheets)
                            if file_data_dict is not None:
                                # 제목 입력 값 가져오기
                                report_html +=  f"<h3>{idx + 1}. {row['제목']}</h3>\n"
                                for sheet_name, df in file_data_dict.items():
                                    wb = openpyxl.load_workbook(file_content)
                                    ws = wb[sheet_name]
                                    html_data = convert_df_to_html_with_styles_and_merging(ws, df)
                                    report_html += f"<div style='text-indent: 20px;'>{html_data}</div>\n"
    
                        else:
                            file_data = extract_data_from_file(file_content, file_type)
                            report_html += f"<h3>{idx + 1}. {row['제목']}</h3>\n<p>{file_data}</p>"
                        if idx > 0 :
                            report_html += "<p/>"
                        html_viewer_data += report_html    
                        #file_data_list.append(row['데이터'])
                    st.session_state['html_report'] = html_viewer_data
                time.sleep(1)  # 예를 들어, 5초 동안 로딩 상태 유지

            with st.spinner('결과 보고서 작성 중입니다...'):
                # LLM 함수 호출
                titles = [row['제목'] for row in st.session_state['rows']]
                requests = [row['요청'] for row in st.session_state['rows']]
        
                responses = run_llm_with_file_and_prompt(
                    st.session_state["openai_api_key"], 
                    titles, 
                    requests, 
                    st.session_state['html_report']
                )
                st.session_state["response"] = responses
                time.sleep(1)  # 예를 들어, 5초 동안 로딩 상태 유지


with col3:
    st.write("")           

# 8 프레임
#st.subheader("")
# 결과 보고서
st.markdown(
    "<p style='font-size:18px; font-weight:bold; color:#007BFF;'>결과 보고서</p>",
    unsafe_allow_html=True
)


# 9 프레임
# LLM 응답 보기
html_result_value = "<div id='html_result_value'>"
with st.expander("📊 결과 보고서 보기", expanded=st.session_state['check_result']):
    if "response" in st.session_state:
        
        for idx, response in enumerate(st.session_state["response"]):
            #st.text_area(f"응답 {idx+1}:", value=response, height=300)
            st.write("결과 보고서 완성")
            html_response_value = f"<div style='border: 0px solid #cccccc; padding: 1px;'>{response}</div>"
            html_result_value += html_response_value
            st.components.v1.html(html_response_value, height=1024, scrolling=True)
    html_result_value += "</div>"
    st.markdown(
        "<hr style='border-top:1px solid #dddddd;border-bottom:0px solid #dddddd;width:100%;padding:0px;margin:0px'></hr>",
        unsafe_allow_html=True
    )
    
# 10 프레임
# 결과 저장 버튼
    col1, col2, col3 = st.columns([0.21, 0.39, 0.4])
    with col1:
        st.markdown(
            "<p style='font-size:14px; font-weight:bold; color:#000000; margin-top:20px;text-align:center;'>결과 보고서<br/>기준일자 지정</p>",
            unsafe_allow_html=True
        )
    with col2:
        # 오늘 날짜 가져오기
        today = datetime.date.today()
        
        # 'report_date_str' 세션 값이 있는지 확인하고, 없으면 'YYYYMMDD' 형식으로 today 값 설정
        if 'report_date_str' not in st.session_state:
            st.session_state['report_date_str'] = today.strftime('%Y%m%d')
        
        
        # 세션에 저장된 'YYYYMMDD' 형식을 date 객체로 변환
        saved_date = today
        # 날짜 문자열을 검사하여 잘못된 형식일 때 예외 처리
        if 'report_date_str' in st.session_state and st.session_state['report_date_str']:
            try:
                # 저장된 날짜 문자열이 있으면 파싱
                saved_date = datetime.datetime.strptime(st.session_state['report_date_str'], '%Y%m%d').date()
            except ValueError:
                # 날짜 형식이 맞지 않으면 오늘 날짜로 설정
                st.warning("잘못된 날짜 형식입니다. 기본값으로 오늘 날짜를 사용합니다.")
        else:
            # 저장된 날짜가 없거나 빈 문자열일 경우 오늘 날짜로 설정
            saved_date = today
    
        report_date = st.date_input(
            "보고서 기준일자 선택",
            value=saved_date,
            min_value=datetime.date(2000, 1, 1),
            max_value=today,
            key="report_date"
        )
        # 날짜를 YYYYMMDD 형식으로 변환
        # 날짜 데이터 메모리에 저장
        st.session_state['report_date_str'] = report_date.strftime("%Y%m%d")
    with col3:
        st.markdown(
            "<p style='font-size:14px; font-weight:bold; color:#000000; margin-top:20px;text-align:left;'>⬅️결과 보고서 저장을 위해 기준일자를 지정해주세요.</p>",
            unsafe_allow_html=True
        )
    col1, col2 = st.columns([0.5, 0.5])
    with col1:   
        if st.button("결과 내용 저장", key="save_result", use_container_width=True):
           
            if "response" in st.session_state:                
                folder_name = st.session_state['selected_folder_name']
                report_date_str = st.session_state.get('report_date_str', datetime.datetime.now().strftime('%Y%m%d'))
                
                # save_html_response 함수를 사용하여 HTML 파일 저장
                file_name, temp_file_path = save_html_response(html_result_value, folder_name, report_date_str)

                # 파일 저장 경로 (reportFiles/{폴더명}/{일자})
                github_folder = f"reportFiles/{folder_name}/{report_date_str}"

                # 폴더 존재 확인 및 생성
                check_and_create_github_folder(github_folder, st.session_state['github_repo'], st.session_state['github_branch'], st.session_state['github_token'])
                
                # GitHub에 HTML 파일 저장
                sha = get_file_sha(st.session_state['github_repo'], f"{github_folder}/{file_name}", st.session_state['github_token'], branch=st.session_state['github_branch'])
                upload_file_to_github(st.session_state['github_repo'], github_folder, file_name, open(temp_file_path, 'rb').read(), st.session_state['github_token'], branch=st.session_state['github_branch'], sha=sha)

                st.success(f"{file_name} 파일이 생성되었습니다.")
                st.download_button(
                    label="다운로드",
                    use_container_width=True,
                    data=open(temp_file_path, 'r', encoding='utf-8').read(),
                    file_name=file_name,
                    mime="text/html"
                )

            else:
                st.warning("결과 보고서를 먼저 실행하세요.")
    with col2:
        if st.button("보고서 양식 저장", key="save_template", use_container_width=True):
            st.session_state['check_result'] = True
            st.session_state['check_report'] = False
            st.session_state['check_upload'] = False
            st.session_state[''] = False
            st.session_state['check_request'] = False
            save_template_to_json()


# 11 프레임
# 결과 보고서 HTML 보기
#if "html_report" in st.session_state:
    #st.write("파일 데이터 추출 보기")
    #html_report_value = f"<div style='border: 2px solid #cccccc; padding: 2px;'>{st.session_state['html_report']}</div>"
    #st.components.v1.html(html_report_value, height=10240, scrolling=True)

# 12 프레임
# 전달된 프롬프트
#st.text_area("전달된 프롬프트:", value="\n\n".join(global_generated_prompt), height=150)
    
# Frontend 기능 구현 끝 ---

