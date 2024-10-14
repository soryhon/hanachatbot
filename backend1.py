import os
import openai
import streamlit as st
from langchain.prompts import PromptTemplate
from langchain import LLMChain
import requests
import PyPDF2
import pandas as pd
import docx
import pptx
from PIL import Image
from io import BytesIO
import base64
import urllib.parse
from openai.error import RateLimitError
from langchain.chat_models import ChatOpenAI
import time
import json
import openpyxl
from openpyxl.utils import get_column_letter
import re
import tempfile
import datetime
import matplotlib.pyplot as plt
import matplotlib.font_manager as fm
import altair as alt
import numpy as np
from langchain.document_loaders import YoutubeLoader
from langchain.chat_models import ChatOpenAI
from langchain.prompts import PromptTemplate
from langchain.chains import LLMChain
from youtube_transcript_api import YouTubeTranscriptApi

# Backend 기능 구현 시작 ---

# 전역변수로 프롬프트 및 파일 데이터 저장
global_generated_prompt = []

# GitHub 정보 및 OpenAI API 키 자동 설정 또는 입력창을 통해 설정
def load_env_info():
    json_data = '''
    {
        "github_repo": "soryhon/hanachatbot",
        "github_branch": "main"
    }
    '''
    
    # JSON 데이터에서 정보 추출
    github_info = json.loads(json_data)
    github_repo = github_info['github_repo']
    github_branch = github_info['github_branch']
    
    # 환경 변수에서 GitHub Token 및 OpenAI API Key 가져오기
    github_token = os.getenv("GITHUB_TOKEN")
    openai_api_key = os.getenv("OPENAI_API_KEY")

    github_set = False
    openai_set = False

    # 입력창을 가로로 배치 (각각 50%의 너비)
    col1, col2 = st.columns(2)

    with col1:
        if not github_token:
            github_token = st.text_input("GitHub Token을 입력하세요", type="password", key="github_token_input")
        else:
            github_set = True
            #st.success("GitHub 정보가 자동으로 설정되었습니다!")
        st.session_state["github_token"] = github_token

    with col2:
        if not openai_api_key:
            openai_api_key = st.text_input("OpenAI API 키를 입력하세요", type="password", key="openai_api_key_input")
        else:
            openai_set = True
            #st.success("OpenAI API 키가 자동으로 설정되었습니다!")
        st.session_state["openai_api_key"] = openai_api_key

    # GitHub 저장소 정보 세션에 저장
    st.session_state["github_repo"] = github_repo
    st.session_state["github_branch"] = github_branch

    # GitHub 정보가 설정되었는지 확인하고 세션 상태 반영
  
    return github_set

# GitHub에 폴더가 존재하는지 확인하고 없으면 폴더를 생성하는 함수
def check_and_create_github_folder_if_not_exists(repo, folder_name, token, branch='main'):
    folder_path = f"{folder_name}/.gitkeep"  # Git에서 빈 폴더 유지용 .gitkeep 파일
    url = f"https://api.github.com/repos/{repo}/contents/{folder_path}?ref={branch}"
    headers = {"Authorization": f"token {token}"}
    
    # 폴더 존재 여부 확인
    response = requests.get(url, headers=headers)
    
    # 폴더가 존재하는 경우
    if response.status_code == 200:
        return True
    
    # 폴더가 없는 경우 (404)
    elif response.status_code == 404:
        #st.warning(f"'{folder_name}' 폴더가 존재하지 않아 생성 중입니다.")
        create_folder_url = f"https://api.github.com/repos/{repo}/contents/{folder_path}"
        data = {
            "message": f"Create {folder_name} folder with .gitkeep",
            "content": base64.b64encode(b'').decode('utf-8'),  # 빈 파일 생성
            "branch": branch
        }
        create_response = requests.put(create_folder_url, json=data, headers=headers)
        
        if create_response.status_code in [200, 201]:
            st.success(f"'{folder_name}' 폴더가 성공적으로 생성되었습니다.")
            return True
        else:
            st.error(f"폴더 생성 실패: {create_response.status_code}")
            return False
    
    # 그 외 상태 코드 처리
    else:
        st.error(f"폴더 확인 중 오류 발생: {response.status_code}")
        return False


# GitHub에서 uploadFiles 하위의 폴더 리스트를 가져오는 함수
def get_folder_list_from_github(repo, branch, token, base_folder='uploadFiles'):

    folder_check = check_and_create_github_folder_if_not_exists(repo, base_folder, token, branch)  # 폴더 생성 및 .gitkeep 파일 추가
    if folder_check:
        url = f"https://api.github.com/repos/{repo}/contents/{base_folder}?ref={branch}"
        headers = {"Authorization": f"token {token}"}
        response = requests.get(url, headers=headers)
        
        # 폴더가 없을 경우 (404) 폴더를 생성하고 다시 폴더 리스트를 가져옴
        if response.status_code == 404:
            st.warning(f"'{base_folder}' 폴더가 존재하지 않아 생성 중입니다.")
            
            response = requests.get(url, headers=headers)  # 폴더 생성 후 다시 요청
            
            if response.status_code != 200:
                st.error("폴더를 생성했지만 다시 폴더 리스트를 가져오지 못했습니다.")
                return []
        
        if response.status_code == 200:
            folders = [item['name'] for item in response.json() if item['type'] == 'dir']
            return folders
        else:
            st.error("폴더 리스트를 가져오지 못했습니다. 저장소 정보나 토큰을 확인하세요.")
            return []
    else:
        return [] 
        
# GitHub에 새로운 폴더를 생성하는 함수
def create_new_folder_in_github(repo, folder_name, token, branch='main'):
    base_folder = "uploadFiles"
    folder_path = f"{base_folder}/{folder_name}/.gitkeep"  # Git에서 빈 폴더를 유지하는 방법 중 하나인 .gitkeep 파일 사용
    url = f"https://api.github.com/repos/{repo}/contents/{folder_path}"
    headers = {"Authorization": f"token {token}"}
    
    data = {
        "message": f"Create folder {folder_name}",
        "content": base64.b64encode(b'').decode('utf-8'),  # 빈 파일로 폴더 생성
        "branch": branch
    }
    
    response = requests.put(url, json=data, headers=headers)
    
    if response.status_code in [200, 201]:
        #st.success(f"'{folder_name}' 폴더가 성공적으로 생성되었습니다.")
        return True
    elif response.status_code == 422:
        st.warning("이미 존재하는 폴더입니다.")
        return False
    else:
        st.error(f"폴더 생성 실패: {response.status_code}")
        return False
      
# 다양한 파일 형식에서 데이터를 추출하는 함수
def extract_data_from_file(file_content, file_type):
    if file_content is None:
        st.error("파일 내용을 가져오지 못했습니다.")
        return None

    if file_type == 'pdf':
        return extract_text_from_pdf(file_content)
    elif file_type == 'csv':
        return extract_text_from_csv(file_content)
    elif file_type == 'docx':
        return extract_text_from_word(file_content)
    elif file_type == 'pptx':
        return extract_text_from_ppt(file_content)
    elif file_type in ['png', 'jpg', 'jpeg']:
        return extract_text_from_image(file_content)
    elif file_type == 'txt':
        return extract_text_from_txt(file_content)
    elif file_type == 'log':
        return extract_text_from_log(file_content)
    elif file_type == 'wav':  # 음성 파일 추가
        return extract_text_from_audio(file_content, file_type)    
    else:
        st.error(f"{file_type} 형식은 지원되지 않습니다.")
        return None

# PDF 파일에서 텍스트 추출
def extract_text_from_pdf(file_content):
    reader = PyPDF2.PdfReader(file_content)
    text = ''
    for page in range(len(reader.pages)):
        text += reader.pages[page].extract_text()
    return text

# CSV 파일에서 텍스트 추출
def extract_text_from_csv(file_content):
    csv_data = pd.read_csv(file_content)
    return csv_data

# 워드 파일에서 텍스트 추출
def extract_text_from_word(file_content):
    doc = docx.Document(file_content)
    return '\n'.join([para.text for para in doc.paragraphs])

# PPT 파일에서 텍스트 추출
def extract_text_from_ppt(file_content):
    presentation = pptx.Presentation(file_content)
    text = ''
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + '\n'
    return text

# 텍스트 파일에서 텍스트 추출
def extract_text_from_txt(file_content):
    try:
        # BytesIO 객체의 내용을 문자열로 변환
        if isinstance(file_content, BytesIO):
            return file_content.read().decode("utf-8")
        elif isinstance(file_content, str):
            return file_content
        else:
            st.error("알 수 없는 파일 형식입니다.")
            return None
    except Exception as e:
        st.error(f"txt 파일에서 텍스트를 추출하는 중 오류가 발생했습니다: {str(e)}")
        return None

# 로그 파일에서 텍스트 추출
def extract_text_from_log(file_content):
    try:
        # BytesIO 객체의 내용을 문자열로 변환
        if isinstance(file_content, BytesIO):
            return file_content.read().decode("utf-8")
        elif isinstance(file_content, str):
            return file_content
        else:
            st.error("알 수 없는 파일 형식입니다.")
            return None
    except Exception as e:
        st.error(f"log 파일에서 텍스트를 추출하는 중 오류가 발생했습니다: {str(e)}")
        return None

    
# 이미지에서 텍스트 추출 (OCR)
def extract_text_from_image(file_content):
    image = Image.open(file_content)
    return "이미지에서 텍스트를 추출하는 기능은 구현되지 않았습니다."


# GitHub에 폴더가 존재하는지 확인하고 없으면 생성하는 함수
def create_github_folder_if_not_exists(repo, folder_name, token, branch='main'):
    url = f"https://api.github.com/repos/{repo}/contents/{folder_name}?ref={branch}"
    headers = {"Authorization": f"token {token}"}
    response = requests.get(url, headers=headers)
    
    if response.status_code == 404:
        st.warning(f"'{folder_name}' 폴더가 존재하지 않아 생성 중입니다.")
        create_folder_url = f"https://api.github.com/repos/{repo}/contents/{folder_name}"
        data = {
            "message": f"Create {folder_name} folder",
            "content": base64.b64encode(b'').decode('utf-8'),  # 빈 파일로 폴더 생성
            "branch": branch
        }
        requests.put(create_folder_url, json=data, headers=headers)
        #st.success(f"'{folder_name}' 폴더가 성공적으로 생성되었습니다.")
    # 폴더가 이미 존재할 경우 메시지를 표시하지 않음

# GitHub API 요청을 처리하는 함수 (파일 목록을 가져옴)
def get_github_files(repo, branch, token):
    # 보고서명 리스트에서 선택한 폴더가 upload_folder에 저장됨
    folder_name = st.session_state.get('upload_folder', 'uploadFiles')
    
    # upload_folder 하위 폴더 내의 파일을 가져옴
    create_github_folder_if_not_exists(repo, folder_name, token, branch)
    url = f"https://api.github.com/repos/{repo}/git/trees/{branch}?recursive=1"
    headers = {"Authorization": f"token {token}"}
    response = requests.get(url, headers=headers)
    
    if response.status_code == 200:
        files = [item['path'] for item in response.json().get('tree', []) if item['type'] == 'blob' and item['path'].startswith(folder_name)]
        return files
    else:
        st.error("GitHub 파일 목록을 가져오지 못했습니다. 저장소 정보나 토큰을 확인하세요.")
        return []

# GitHub에서 파일의 SHA 값을 가져오는 함수
def get_file_sha(repo, file_path, token, branch='main'):
    encoded_file_path = urllib.parse.quote(file_path)
    url = f"https://api.github.com/repos/{repo}/contents/{encoded_file_path}?ref={branch}"
    headers = {"Authorization": f"token {token}"}
    response = requests.get(url, headers=headers)
    
    if response.status_code == 200:
        return response.json().get('sha', None)
    else:
        return None

# GitHub에 파일 업로드 함수 (덮어쓰기 포함)
def upload_file_to_github(repo, folder_name, file_name, file_content, token, branch='main', sha=None):
    create_github_folder_if_not_exists(repo, folder_name, token, branch)
    encoded_file_name = urllib.parse.quote(file_name)
    url = f"https://api.github.com/repos/{repo}/contents/{folder_name}/{encoded_file_name}"
    headers = {
        "Authorization": f"token {token}",
        "Content-Type": "application/json"
    }

    content_encoded = base64.b64encode(file_content).decode('utf-8')

    data = {
        "message": f"Upload {file_name}",
        "content": content_encoded,
        "branch": branch
    }

    if sha:
        data["sha"] = sha

    response = requests.put(url, json=data, headers=headers)

    if response.status_code in [200, 201]:
        st.success(f"{file_name} 파일이 성공적으로 업로드(덮어쓰기) 되었습니다.")
    else:
        st.error(f"파일 업로드에 실패했습니다: {response.status_code}")
        st.error(response.json())

# GitHub에서 파일을 다운로드하는 함수
def get_file_from_github(repo, branch, filepath, token):
    encoded_filepath = urllib.parse.quote(filepath)
    url = f"https://api.github.com/repos/{repo}/contents/{encoded_filepath}?ref={branch}"
    headers = {"Authorization": f"token {token}"}
    response = requests.get(url, headers=headers)
    
    if response.status_code == 200:
        return BytesIO(requests.get(response.json()['download_url']).content)
    else:
        st.error(f"{filepath} 파일을 가져오지 못했습니다. 상태 코드: {response.status_code}")
        return None

# 엑셀 시트에서 셀 스타일 및 병합 정보 추출
def extract_cell_style_and_merged(ws):
    style_dict = {}
    merged_cells = ws.merged_cells.ranges
    
    for row in ws.iter_rows():
        for cell in row:
            cell_style = {
                "alignment": cell.alignment.horizontal if cell.alignment else 'left',
                "font_bold": cell.font.bold if cell.font else False,
                "border": bool(cell.border and (cell.border.left or cell.border.right or cell.border.top or cell.border.bottom)),
                "value": "" if str(cell.value).startswith("Unnamed") else cell.value
            }
            style_dict[cell.coordinate] = cell_style
            
    return style_dict, merged_cells

# 엑셀 시트 데이터를 HTML로 변환하고 스타일 및 병합 적용
def convert_df_to_html_with_styles_and_merging(ws, df):
    # 1행의 'Unnamed: 숫자' 형식 값은 공백으로 처리
    df.columns = [re.sub(r'Unnamed: \d+', '', str(col)).strip() for col in df.columns]
    
    style_dict, merged_cells = extract_cell_style_and_merged(ws)
    df = df.fillna('')  # NaN 값을 공백으로 처리
    html = "<table class='table table-bordered' style='border-spacing: 0; border-collapse: collapse;'>\n"

    # 헤더 부분 (border 간격 없게 설정)
    html += "<thead>\n<tr>\n"
    for col in df.columns:
        html += f"<th style='text-align:center; font-weight:bold; background-color:#E7E6E6; border: 1px solid black;'>{col}</th>\n"
    html += "</tr>\n</thead>\n"

    # 병합 셀 정보 저장
    merged_cells_dict = {}
    for merged in merged_cells:
        min_row, min_col, max_row, max_col = merged.min_row, merged.min_col, merged.max_row, merged.max_col
        for row in range(min_row, max_row + 1):
            for col in range(min_col, max_col + 1):
                merged_cells_dict[(row, col)] = (min_row, min_col, max_row, max_col)

    # 데이터 부분
    html += "<tbody>\n"
    for row_idx, row in df.iterrows():
        html += "<tr>\n"
        for col_idx, value in enumerate(row):
            cell_ref = f"{get_column_letter(col_idx+1)}{row_idx+2}"  # 셀 참조 계산
            style = style_dict.get(cell_ref, {})
            alignment = style.get("alignment", "left")
            font_weight = "bold" if style.get("font_bold", False) else "normal"
            border = "1px solid black" if style.get("border", False) else "none"
            cell_coordinates = (row_idx + 2, col_idx + 1)

            if cell_coordinates in merged_cells_dict:
                min_row, min_col, max_row, max_col = merged_cells_dict[cell_coordinates]
                rowspan = max_row - min_row + 1
                colspan = max_col - min_col + 1
                if (row_idx + 2, col_idx + 1) == (min_row, min_col):
                    html += f"<td rowspan='{rowspan}' colspan='{colspan}' style='text-align:{alignment}; font-weight:{font_weight}; border:{border};'>{value}</td>\n"
            elif cell_coordinates not in merged_cells_dict:
                html += f"<td style='text-align:{alignment}; font-weight:{font_weight}; border:{border};'>{value}</td>\n"
        html += "</tr>\n"
    html += "</tbody>\n</table>"
    
    return html

# 엑셀 파일에서 시트 가져오기 및 데이터 추출
def extract_sheets_from_excel(file_content, selected_sheets):
    try:
        excel_data = pd.ExcelFile(file_content)
        all_sheets = excel_data.sheet_names
        
        # 사용자가 선택한 시트를 가져옴
        if selected_sheets == 'all':
            selected_sheets = all_sheets
        else:
            selected_sheets = [all_sheets[int(i)-1] for i in selected_sheets if int(i) <= len(all_sheets)]
        
        # 선택한 시트의 데이터를 하나로 합침
        data_dict = {}
        with pd.ExcelFile(file_content) as xls:
            for sheet in selected_sheets:
                df = pd.read_excel(xls, sheet_name=sheet)
                data_dict[sheet] = df
                
        return data_dict
    except Exception as e:
        st.error(f"엑셀 파일의 시트 데이터를 추출하는 중에 오류가 발생했습니다: {str(e)}")
        return None

# 숫자, ',' 및 '-'만 허용하는 함수
def validate_sheet_input(input_value):
    if all(c.isdigit() or c in ['-', ','] for c in input_value):
        return True
    return False

# 시트 선택 로직 추가
def handle_sheet_selection(file_content, sheet_count, idx):
    # 3개의 객체를 가로로 배치
    col1, col2 = st.columns([0.5, 0.5])
    with col1:
        st.text_input(f"시트 갯수_{idx}", value=f"{sheet_count}개", key=f"sheet_count_{idx}", disabled=True)  # 시트 갯수 표시 (비활성화)
    
    with col2:
        # st.session_state['rows']와 st.session_state['rows'][idx]['파일정보']가 유효한지 확인하여 값 설정
        sheet_selection_value = '1'
        if st.session_state.get('rows') and st.session_state['rows'][idx].get('파일') and st.session_state['rows'][idx].get('파일정보'):
            sheet_selection_value = st.session_state['rows'][idx]['파일정보']
            file_name = st.session_state['rows'][idx]['파일']
        
        sheet_selection = st.text_input(
            f"시트 선택_{idx}(예: 1-3, 5)", 
            value=sheet_selection_value, 
            key=f"sheet_selection_{idx}_{file_name}"
        )
    
        # 입력값 변경 시 세션에 저장
        if validate_sheet_input(sheet_selection):
            st.session_state['rows'][idx]['파일정보'] = sheet_selection
        else:
            st.error("잘못된 입력입니다. 숫자와 '-', ',' 만 입력할 수 있습니다.")
    return None

# 시트 선택 입력값을 분석하는 함수
def parse_sheet_selection(selection, sheet_count):
    selected_sheets = []

    try:
        if '-' in selection:
            start, end = map(int, selection.split('-'))
            if start <= end <= sheet_count:
                selected_sheets.extend(list(range(start, end+1)))
        elif ',' in selection:
            selected_sheets = [int(i) for i in selection.split(',') if 1 <= int(i) <= sheet_count]
        else:
            selected_sheets = [int(selection)] if 1 <= int(selection) <= sheet_count else []
    except ValueError:
        st.error("잘못된 시트 선택 입력입니다.")
        return None

    return selected_sheets

# 파일에서 데이터를 추출하고 요청사항 리스트에서 선택한 엑셀 파일의 시트를 보여주는 로직 수정
def handle_file_selection(file_path, file_content, file_type, idx):
    if file_type == 'xlsx':
        wb = openpyxl.load_workbook(file_content)
        sheet_count = len(wb.sheetnames)

        # 시트 선택 로직 처리
        handle_sheet_selection(file_content, sheet_count, idx)
        #return file_data_dict
    #else:
        #return extract_data_from_file(file_content, file_type)

# HTML 보고서 생성 함수 (배열에서 데이터 가져옴)
def generate_final_html_report(file_data):
    report_html = ""
    if "html_report" in st.session_state:
        report_html = st.session_state['html_report']
    if file_data:
        report_html += f"<div style='text-indent: 1px;'>\n{file_data}\n</div><p/.\n"
        st.session_state['html_report'] = report_html  # 최종 값을 세션 상태에 저장

# 엑셀 데이터 및 을 HTML로 변환하여 하나의 세트로 출력하는 함수
def generate_html_report_with_title(titles, data_dicts):
    report_html = ""
    
    for i, (title, data_dict) in enumerate(zip(titles, data_dicts), start=1):
        report_html += f"<h3>{i}. {title}</h3>\n"
        report_html += "<div style='text-indent: 20px;'>\n"
        
        for sheet_name, df in data_dict.items():
            wb = openpyxl.load_workbook(BytesIO(df))
            ws = wb[sheet_name]
            report_html += convert_df_to_html_with_styles_and_merging(ws, df)
        
        report_html += "</div>\n"
    
    return report_html

# LLM 연동 및 응답 처리 함수
def execute_llm_request(api_key, prompt):
    openai.api_key = api_key
    responses = []

    try:
        # 텍스트 길이 제한 확인 (예: 4000자로 제한)
        if len(prompt) > 4000:
            st.error("프롬프트 글자 수 초과로 LLM 연동에 실패했습니다.")
            return None

        # 프롬프트 템플릿 설정
        prompt_template = PromptTemplate(
            template=prompt,
            input_variables=[]
        )

        # LLM 모델 생성
        llm = ChatOpenAI(model_name="gpt-4o")
        chain = LLMChain(llm=llm, prompt=prompt_template)

        success = False
        retry_count = 0
        max_retries = 5  # 최대 재시도 횟수

        # 응답을 받을 때까지 재시도
        while not success and retry_count < max_retries:
            try:
                response = chain.run({})
                responses.append(response)
                success = True
            except RateLimitError:
                retry_count += 1
                st.warning(f"API 요청 한도를 초과했습니다. 10초 후 다시 시도합니다. 재시도 횟수: {retry_count}/{max_retries}")
                time.sleep(10)

            time.sleep(10)
    except Exception as e:
        st.error(f"LLM 실행 중 오류가 발생했습니다: {str(e)}")

    return responses

# LLM을 통해 프롬프트와 파일을 전달하고 응답을 받는 함수
def run_llm_with_file_and_prompt(api_key, titles, requests, file_data_str):
    global global_generated_prompt
    openai.api_key = api_key

    responses = []
    global_generated_prompt = []  # 프롬프트들을 담을 리스트 초기화

    try:
        # 요청사항 리스트 문자열 생성
        request_list_str = "\n".join([
            f"{i+1}.{title}의 항목 데이터에 대해 '{request}' 요청 사항을 만족하게 구성한다. 항목 데이터의 데이터 값은 중략하거나 누락되어서 안된다.\n"
            for i, (title, request) in enumerate(zip(titles, requests))
        ])

        # 프롬프트 텍스트 정의
        generated_prompt = f"""
        아애의 항목 데이터를 간결하고 깔끔한 보고서 작성을 위해 보고서 내용에 대해서 알기 쉽게 내용 요약하고 설명해야 한다.
        다음과 같은 조건에 모두 만족해야 한다.
        가. 아래의 항목 데이터를 분석하여 각 항목마다의 '요청사항' 리스트와 조건사항에 대해 모두 만족할 수 있도록 최적화된 보고서를 완성해.
        나. 항목 데이터 내 가장 첫번째 행은 각 보고서 항목에 타이틀이므로 순번과 문구를 그대로 유지해야 한다. 이 항목의 타이틀을 기준으로 각 항목 데이터를 분류하고 그에 맞는 요청사항을 반영해야 한다.
        다. 항목데이터 중 table 태그가 포함된 데이터는 엑셀에서 추출한 데이터로 표 형식으로 그대로 유지해야 한다. 
              table태그 데이터의 데이터 값은 중략하거나 누락되어서 안된다.보고서에 꼭 필요하니 최대한 그대로 구조와 데이터는 출력되게 하고 내용만 업데이트한다. 
        라. 표 형식의로 table태그로 답변 할 때는 th과 td 태그는 border는 사이즈 1이고 색상은 검정색으로 구성한다. table 태그 가로길이는 전체를 차지해야 한다.
        마. 표 형식으로 답변할 때는 반드시 모든 항목 데이터의 수정한 데이터 내용과 HTML 태그를 보완한 모든 데이터를 업그레이드한 데이터로 보여줘야 한다.
        바. 이외 table 태그가 포함 안된 데이터는 파일에서 텍스트를 추출한 데이터로 내용으로 보고서 양식에 맞춰 간결하고 깔끔하게 요약하고 html 형식으로 변환해야 한다.
        사. 답변할 때는 반드시 모든 항목 데이터의 수정한 데이터 내용과 HTML 형삭에 맞춰 답변한다. 문단마다 줄바꿈을 적용하여 br태그 활용하고 가시성 높게 특수기호를 활용하여 보고서 양식에 준하게 요약한 내용을 설명해줘야 한다.
        아. 항목 데이터를 업데이트 하는 결과 가장 먼저 나와야하고 위에는 그 어떤한 설명 내용이 응답하면 안 된다. 
        자. 그 다음 일정 간격을 두고 h3 태그를 활용해서 '✨AI 요약과 설명' 타이틀 추가하고 색상을 달리 구성한다. 너의 답변이라는 것을 표현하는 특수문자로 강조해.
               전달받은 보고서 전반적인 내용에 대해 너가 선정한 가장 좋은 방법으로 요약과 설명하고 그 내용을 HTML 형식으로 변환하여 답변해야 한다.
        차. '````', '````HTML' 이 문구들이 답변에 포함되지 않아야 한다.
        -요청사항
        [
            {request_list_str}
        ]
        -항목 데이터
        [
            {file_data_str}
        ]
        """
        
        # 텍스트 길이 제한 확인 (예: 1000000자로 제한)
        if len(generated_prompt) > 1000000:
            st.error("프롬프트 글자 수 초과로 LLM 연동에 실패했습니다.")
        else:
            global_generated_prompt.append(generated_prompt)
            prompt_template = PromptTemplate(
                template=generated_prompt,
                input_variables=[]
            )

            # LLM 모델 생성
            llm = ChatOpenAI(model_name="gpt-4o")
            chain = LLMChain(llm=llm, prompt=prompt_template)

            success = False
            retry_count = 0
            max_retries = 5  # 최대 재시도 횟수

            # 응답을 받을 때까지 재시도
            while not success and retry_count < max_retries:
                try:
                    response = chain.run({})
                    responses.append(response)
                    success = True
                except RateLimitError:
                    retry_count += 1
                    st.warning(f"API 요청 한도를 초과했습니다. 10초 후 다시 시도합니다. 재시도 횟수: {retry_count}/{max_retries}")
                    time.sleep(10)

                time.sleep(10)
    except Exception as e:
        st.error(f"LLM 실행 중 오류가 발생했습니다: {str(e)}")

    return responses


#st.session_state를 새로고침하는 함수
def refresh_page():
    if 'is_updating' not in st.session_state:
        st.session_state['is_updating'] = False
    elif st.session_state['is_updating']:
        st.session_state['is_updating'] = False
    else:
        st.session_state['is_updating'] = True

def init_session_state(check_value):
    folderlist_init_value = "보고서명을 선택하세요."
    templatelist_init_value = "불러올 보고서 양식을 선택하세요."
    if(check_value):
            st.session_state['rows'] = [
                {"제목": "", "요청": "", "파일": "", "데이터": "","파일정보":"1" }
                for _ in range(st.session_state['num_requests'])
            ]    
            st.session_state['html_report'] = ""
    else:
        if 'selected_folder_name' not in st.session_state:
            st.session_state['selected_folder_name'] = folderlist_init_value
        if 'folder_list_option' not in st.session_state:       
            st.session_state['folder_list_option'] = folderlist_init_value
        if 'selected_template_name' not in st.session_state:
            st.session_state['selected_template_name'] = templatelist_init_value
        if 'template_list_option' not in st.session_state:       
            st.session_state['template_list_option'] = templatelist_init_value
        if 'upload_folder' not in st.session_state:        
            st.session_state['upload_folder'] = "uploadFiles" 
        if 'selected_folder_index' not in st.session_state:    
            st.session_state['selected_folder_index'] = 0
        if 'selected_template_index' not in st.session_state:
            st.session_state['selected_template_index'] = 0
        if 'new_folder_text' not in st.session_state:    
            st.session_state['new_folder_text'] = ""
        if 'check_report' not in st.session_state:    
            st.session_state['check_report'] = True
        if 'check_upload' not in st.session_state:    
            st.session_state['check_upload'] = False        
        if 'check_request' not in st.session_state:    
            st.session_state['check_request'] = False
        if 'check_result' not in st.session_state:    
            st.session_state['check_result'] = False
        if 'check_setting' not in st.session_state:    
            st.session_state['check_setting'] = False
        if 'report_date_str' not in st.session_state: 
            st.session_state['report_date_str'] = ""
# HTML 파일을 저장하고 파일 경로를 반환하는 함수 (날짜 포함)
def save_html_response(html_content, folder_name, report_date_str):
    # 현재 시간을 'YYYYMMDDHHMMSS' 형식으로 가져오기
    #current_time = datetime.datetime.now().strftime("%Y%m%d%H%M%S")
    # HTML 파일명을 보고서명과 날짜로 설정
    file_name = f"{folder_name}_result_{report_date_str}.html"
    
    # HTML 파일 임시 경로에 저장
    temp_file_path = f"/tmp/{file_name}"
    with open(temp_file_path, 'w', encoding='utf-8') as f:
        f.write(html_content)
    
    return file_name, temp_file_path        

# GitHub에 폴더가 존재하는지 확인하는 함수
def check_and_create_github_folder(folder_name, repo, branch, token):
    url = f"https://api.github.com/repos/{repo}/contents/{folder_name}"
    headers = {"Authorization": f"token {token}"}

    # 폴더 존재 여부 확인
    response = requests.get(url, headers=headers)
    if response.status_code == 404:  # 폴더가 없으면 생성
        # 폴더 생성하기 위한 커밋 메시지와 파일 내용 설정
        data = {
            "message": f"Create {folder_name} folder",
            "content": "",  # GitHub에서는 폴더 자체를 직접 생성할 수 없으므로 빈 파일 생성으로 대체
            "branch": branch
        }
        # 빈 파일 생성 (ex: .gitkeep)
        file_url = f"{url}/.gitkeep"
        response = requests.put(file_url, headers=headers, data=json.dumps(data))
        if response.status_code == 201:
            st.success(f"{folder_name} 폴더가 생성되었습니다.")
        else:
            st.error(f"{folder_name} 폴더 생성 실패: {response.json()}")
    elif response.status_code == 200:
        #st.info(f"{folder_name} 폴더가 이미 존재합니다.")
        return None
    else:
        st.error(f"폴더 확인 중 오류 발생: {response.json()}")
        
# JSON 파일 저장 함수
def save_template_to_json():
    repo = st.session_state["github_repo"]
    branch = st.session_state["github_branch"]
    token = st.session_state["github_token"]

    # GitHub 토큰과 레포지토리 설정 확인
    if not token or not repo:
        st.error("GitHub 토큰이나 저장소 정보가 설정되지 않았습니다.")
        return
        
    # JSON 데이터 구조 생성
    template_data = {
        "selected_folder_name": st.session_state['selected_folder_name'],
        "num_requests": st.session_state['num_requests'],
        "rows": st.session_state['rows'],
        "rows_length": len(st.session_state['rows']),
        "timestamp": datetime.datetime.now().strftime("%Y%m%d%H%M%S")
    }

    # 파일명 생성
    folder_name = st.session_state['selected_folder_name']
    timestamp = template_data["timestamp"]
    json_file_name = f"{folder_name}_Template_{timestamp}.json"

    # GitHub 저장소 내 templateFiles 폴더 생성 및 파일 저장
    template_folder = "templateFiles"
    check_and_create_github_folder(template_folder, repo, branch, token)
   
    # 저장할 파일 경로
    json_file_path = f"{template_folder}/{json_file_name}"

    # JSON 파일을 Base64로 인코딩
    json_content = json.dumps(template_data, ensure_ascii=False, indent=4)
    json_base64 = base64.b64encode(json_content.encode('utf-8')).decode('utf-8')

    # GitHub에 파일 업로드
    url = f"https://api.github.com/repos/{repo}/contents/{json_file_path}"
    headers = {"Authorization": f"token {token}"}
    data = {
        "message": f"Add {json_file_name}",
        "content": json_base64,
        "branch": branch
    }

    response = requests.put(url, headers=headers, json=data)
    if response.status_code == 201:
        st.success(f"{json_file_name} 파일이 {template_folder} 폴더에 저장되었습니다.")
    else:
        st.error(f"파일 저장 실패: {response.json()}")

# GitHub에서 templateFiles 폴더 내의 JSON 파일 리스트를 가져오는 함수
def get_template_files_list(repo, branch, token):
    template_folder = "templateFiles"
    url = f"https://api.github.com/repos/{repo}/contents/{template_folder}?ref={branch}"
    headers = {"Authorization": f"token {token}"}
    response = requests.get(url, headers=headers)

    if response.status_code == 200:
        # JSON 파일만 필터링하여 리스트로 반환
        return [item['name'] for item in response.json() if item['name'].endswith('.json')]
    else:
        st.error("templateFiles 폴더의 파일 목록을 가져오지 못했습니다.")
        return []

# JSON 파일의 내용을 불러오는 함수
def load_template_from_github(repo, branch, token, file_name):
    template_folder = "templateFiles"
    json_file_path = f"{template_folder}/{file_name}"
    url = f"https://api.github.com/repos/{repo}/contents/{json_file_path}?ref={branch}"
    headers = {"Authorization": f"token {token}"}
    response = requests.get(url, headers=headers)

    if response.status_code == 200:
        file_content = base64.b64decode(response.json()['content'])
        return json.loads(file_content)
    else:
        st.error(f"{file_name} 파일을 가져오지 못했습니다.")
        return None

def apply_template_to_session_state(file_name):
    try:
        # 템플릿 JSON 파일 로드
        with open(file_name, 'r', encoding='utf-8') as f:
            template_data = json.load(f)
        
        # JSON 데이터에서 세션 상태 적용
        selected_folder_name = template_data.get('selected_folder_name', '')
        #num_requests = template_data.get('num_requests', 1)
        rows = template_data.get('rows', [])
        
        # 세션 상태에 값 저장

        st.session_state['selected_folder_name'] = selected_folder_name
        st.session_state['rows'] = rows
        st.session_state['is_updating'] = False
        st.session_state['upload_folder'] = f"uploadFiles/{selected_folder_name}"
        st.session_state['check_report'] = False
        st.session_state['check_upload'] = False
        st.session_state['check_request'] = True
        st.session_state['check_result'] = False
        st.session_state['selected_folder_index'] = 0
    
        # 'num_requests'는 직접 변경할 수 없으므로 Streamlit에서 제공하는 방법으로 값을 설정
        #if "num_requests" in st.session_state:
            #st.session_state["num_requests"] = num_requests

        
        # folder_list에서 selected_folder_name의 인덱스 찾기
        folder_list = st.session_state.get('folder_list_option', [])
        if selected_folder_name in folder_list:
            selected_index = folder_list.index(selected_folder_name)
            st.session_state['selected_folder_index'] = selected_index + 1
        
        # 엑셀 파일 처리: 파일 정보에 따라 시트 선택 입력창 추가
        for idx, row in enumerate(rows):
            selected_file_name = row.get("파일")
            file_info = row.get("파일정보", "1")
            
            if selected_file_name and selected_file_name.endswith('.xlsx'):
                # 시트 선택 로직 적용
                #file_content = get_file_from_github(
                    #st.session_state["github_repo"],
                    #st.session_state["github_branch"],
                    #selected_file_name,
                    #st.session_state["github_token"]
                #)
                #if file_content:
                    #handle_sheet_selection(file_content, len(openpyxl.load_workbook(file_content).sheetnames), idx)
                st.session_state['rows'][idx]['파일정보'] = file_info        
        st.success(f"'{selected_folder_name}' 양식을 불러오기 성공하였습니다.")
    
    except FileNotFoundError:
        st.error(f"파일 '{file_name}'을 찾을 수 없습니다.")
    except json.JSONDecodeError:
        st.error(f"'{file_name}' 파일을 파싱하는 중 오류가 발생했습니다. JSON 형식을 확인해주세요.")
    except Exception as e:
        st.error(f"템플릿 불러오기 중 오류가 발생했습니다: {str(e)}")

# 보고서명 리스트를 가져오고, reportFiles 폴더 존재 여부를 확인하고 없으면 생성하는 함수
def get_report_folder_list_from_github(repo, branch, token):
    base_folder = "reportFiles"
    folder_check = check_and_create_github_folder_if_not_exists(repo, base_folder, token, branch)

    if folder_check:
        url = f"https://api.github.com/repos/{repo}/contents/{base_folder}?ref={branch}"
        headers = {"Authorization": f"token {token}"}
        response = requests.get(url, headers=headers)

        if response.status_code == 200:
            folder_list = [item['name'] for item in response.json() if item['type'] == 'dir']
            return folder_list
        else:
            st.error(f"폴더 리스트를 가져오지 못했습니다: {response.status_code}")
            return []
    else:
        st.error("reportFiles 폴더가 존재하지 않아 생성할 수 없습니다.")
        return []

# 폴더명 리스트와 하위 폴더 리스트 가져오기
def get_subfolder_list(repo, branch, token, selected_folder):
    base_folder = f"reportFiles/{selected_folder}"
    url = f"https://api.github.com/repos/{repo}/contents/{base_folder}?ref={branch}"
    headers = {"Authorization": f"token {token}"}
    response = requests.get(url, headers=headers)

    if response.status_code == 200:
        subfolder_list = [item['name'] for item in response.json() if item['type'] == 'dir']
        date_list = [datetime.datetime.strptime(folder, '%Y%m%d').date() for folder in subfolder_list]
        date_list.sort()  # 날짜 오름차순 정렬
        return subfolder_list, date_list
    else:
        st.error(f"하위 폴더 리스트를 가져오지 못했습니다: {response.status_code}")
        return [], []

# 특정 폴더 내에서 HTML 파일 목록을 가져오는 함수
def get_html_files_from_folder(repo, branch, folder_path, token):
    url = f"https://api.github.com/repos/{repo}/contents/{folder_path}?ref={branch}"
    headers = {"Authorization": f"token {token}"}
    response = requests.get(url, headers=headers)

    if response.status_code == 200:
        # HTML 파일만 필터링하여 리스트로 반환
        return [item['name'] for item in response.json() if item['name'].endswith('.html')]
    else:
        st.error(f"{folder_path} 폴더에서 HTML 파일을 가져오지 못했습니다: {response.status_code}")
        return []

# 시작일자와 마지막 일자를 설정하고 버튼 클릭 시 데이터를 가져오는 함수
def fetch_report_data_between_dates(repo, branch, token, selected_folder, start_date, end_date):
    subfolder_list, date_list = get_subfolder_list(repo, branch, token, selected_folder)

    # 시작일자, 마지막 일자 인덱스 추출
    start_index = max(0, min(range(len(date_list)), key=lambda i: abs(date_list[i] - start_date)))
    end_index = min(range(len(date_list)), key=lambda i: abs(date_list[i] - end_date))

    # 조건에 맞는 폴더들의 데이터를 가져오기
    report_html = ""
    num =0
    for idx in range(start_index, end_index + 1):
        folder_name = subfolder_list[idx]
        num += 1
        report_html += f"<h3>[보고서{num}]-기준일자: {date_list[idx].strftime('%Y년 %m월 %d일')}</h3>\n"
        folder_path = f"reportFiles/{selected_folder}/{folder_name}"

        # 해당 폴더 내 HTML 파일 목록을 가져옴
        html_files = get_html_files_from_folder(repo, branch, folder_path, token)

        # 폴더 내 HTML 파일이 있을 경우, 그 중 첫 번째 파일을 사용
        if html_files:
            html_file_path = f"{folder_path}/{html_files[0]}"
            file_content = get_file_from_github(repo, branch, html_file_path, token)
            if file_content:
                report_html += file_content.read().decode('utf-8')
            else:
                st.error(f"{html_file_path} 파일을 GitHub에서 가져오지 못했습니다.")
        else:
            st.warning(f"{folder_path} 폴더 내에 HTML 파일이 없습니다.")
    return report_html

# LLM을 통해 프롬프트와 파일을 전달하고 응답을 받는 함수
def run_llm_with_analysisfile_and_prompt(api_key, title, request, file_data_str):
    global global_generated_prompt
    openai.api_key = api_key

    responses = []
    global_generated_prompt = []  # 프롬프트들을 담을 리스트 초기화

    try:

        # 프롬프트 텍스트 정의
        generated_prompt = f"""
        아래의 항목 데이터는 여러 보고서 데이터를 하나로 취합한 것으로 '기준일자' 포함 타이틀로 데이터를 분류하고 기준일자별로 변화 추이를 분석하여 설명해.
        다음과 같은 조건에 모두 만족해야 한다.
        가. '기준일자' 포함 타이틀로 데이터를 분류하여 총 몇 개의 보고서 데이터인지 판단해야 하고, 기준일자 통해서 각 분류한 데이터에 일자를 판별해야 한다.
        나. 분류된 각 데이터 내 'AI 요약과 설명'이 포함한 행부터 마지막 행까지는 분석 대상이 아니므로 제외한다.
        다. 이와 같이 데이터를 분류하여 기준일자 기준으로 처음부터 끝까지 보고서를 차례대로 비교하고 분석해야 한다. 여러 데이터중 모두 일치 항목에 대해 선정하고 이 항목에 대해 기준일자별로 데이터 값을 비교할 수 있게 변화추이 표로 완성해야 한다.
        라. 표 형식의로 table태그로 답변 할 때는 th과 td 태그는 border는 사이즈 1이고 색상은 검정색으로 구성한다. table 태그 가로길이는 전체를 차지해야 한다.
        마. 이외 table 태그가 포함 안된 설명은 너가 생각한 가장 좋은 보고서 양식에 맞춰 간결하고 깔끔하게 요약하고 html 형식으로 변환해야 한다.
        바. 답변할 때는 반드시 모든 항목 데이터의 수정한 데이터 내용과 HTML 형삭에 맞춰 답변한다. 문단마다 줄바꿈을 적용하여 br태그 활용하고 가시성 높게 특수기호와 이모지를 활용하여 보고서 양식에 준하게 요약한 내용을 설명한다.
        사. 답변할 때 첫번째 행에는 h3 태그를 활용해서 '{title}' 문구가 반드시 시작되어야 한다.
        아. 답변할 때 두번째 행에는 h3 태그를 활용해서 '✨AI 비교 분석 결과' 타이틀 추가하고 색상을 달리 구성한다. 너의 답변이라는 것을 표현하는 특수문자로 강조해.
               전달받은 보고서 전반적인 내용에 대해 너가 선정한 가장 좋은 방법으로 요약과 설명하고 그 내용을 HTML 형식으로 변환하여 답변해야 한다.
        자. '````', '````HTML' 이 문구들이 답변에 포함되지 않아야 한다.
        차. 'AI 비교 분석 결과'로 비교 분석 설명하고 아래에는 h3 태그를 활용해서 '결과 차트 추천' 타이틀 추가하고 각 데이별 변화된 추이를 차트로 표현할 수 있게 아래의 양식에 맞춰서 답변한다.
             차트는 반드시 python 기반으로 Streamlit에서 구현이 가능한 차트에서 추천해. 답변 양식의 각 항목별 타이틀을 뚜렷하게 강조해야 된다.
        카. 차트 위한 import streamlit as st, import pandas as pd, import matplotlib.pyplot as plt, import altair as alt, import matplotlib.font_manager as fm, import numpy as np 이 선언되어 있으니 import 할 필요 없이 구현하고,
            이외 추가로 필요한 경우만 import 해야한다. 차트 색상은 비비드 컬러에화려한 색상으로 뚜렷하게 하고 가시성 높게 글자 크기도 중간 크기로 설정해야 한다.
        카. 답변 양식 정의
        [
            전체 타이틀
            ✨AI 비교 분석 결과
                1. 분석한 보고서 수 : << !갯수 표기 (예: 2개)
                2. 보고서 설명 : << !줄바꿈은 <br/>로 표기.
                3. 분석한 내용 : << !줄바꿈은 <br/>로 표기. 보고서 양식
                4. 변화추이 표 : << !table 태그로 표현
                5. 차트 설명 : << !줄바꿈은 <br/>로 표기
                [[ ]]         << ![[ ]]사이에 차트 코드만 표기
                
        ]
        -요청사항
        [
            {request}
        ]
        -항목 데이터
        [
            {file_data_str}
        ]
        """
        
        # 텍스트 길이 제한 확인 (예: 1000000자로 제한)
        if len(generated_prompt) > 1000000:
            st.error("프롬프트 글자 수 초과로 LLM 연동에 실패했습니다.")
        else:
            global_generated_prompt.append(generated_prompt)
            prompt_template = PromptTemplate(
                template=generated_prompt,
                input_variables=[]
            )

            # LLM 모델 생성
            llm = ChatOpenAI(model_name="gpt-4o")
            chain = LLMChain(llm=llm, prompt=prompt_template)

            success = False
            retry_count = 0
            max_retries = 5  # 최대 재시도 횟수

            # 응답을 받을 때까지 재시도
            while not success and retry_count < max_retries:
                try:
                    response = chain.run({})
                    response = response.replace('/n', '<br/>')
                    response = response.replace('```html', '')
                    response = response.replace('```', '')
                    responses.append(response)
                    success = True
                except RateLimitError:
                    retry_count += 1
                    st.warning(f"API 요청 한도를 초과했습니다. 10초 후 다시 시도합니다. 재시도 횟수: {retry_count}/{max_retries}")
                    time.sleep(10)

                time.sleep(10)
    except Exception as e:
        st.error(f"LLM 실행 중 오류가 발생했습니다: {str(e)}")

    return responses

def extract_text_within_brackets(response):
    comment_text = ""
    extracted_text = ""
    # responses는 리스트로, 각 응답을 반복하며 정규식으로 텍스트 추출
    if len(response) > 0 :
        start_index = response.find('[[')  # '[[-'의 첫 번째 인덱스
        end_index = response.find(']]', start_index)  # '-]]'의 첫 번째 인덱스
        
        # '[['과 ']]'가 모두 존재할 때만 추출
        if start_index != -1 and end_index != -1 and (start_index < end_index):
            # '[[', ']]'을 포함한 부분을 추출
            extracted_text = response[start_index+2:end_index ]  # ']]'도 포함시키기 위해 +2
            extracted_text = extracted_text.replace('<br/>', '')
            extracted_text = extracted_text.replace('```python', '')
            extracted_text = extracted_text.replace('python', '')
            extracted_text = extracted_text.replace('```', '')
        if start_index != -1 :
            comment_text = response[0:start_index ] 

    return comment_text, extracted_text

# GitHub 저장소에서 폰트 파일을 다운로드하는 함수
def download_and_apply_font_from_github(github_repo, branch, fm, plt):
    font_file_path = "font/NanumGothic.ttf"
    save_dir="/tmp"
    font_url = f"https://raw.githubusercontent.com/{github_repo}/{branch}/{font_file_path}"
    font_name = os.path.basename(font_file_path)  # 파일명 추출
    font_path = os.path.join(save_dir, font_name)  # 저장할 경로

    # 폰트 파일 다운로드
    if not os.path.exists(font_path):
        response = requests.get(font_url)
        if response.status_code == 200:
            with open(font_path, 'wb') as f:
                f.write(response.content)
            #st.success(f"폰트 {font_name} 다운로드 성공")
        else:
            raise Exception(f"폰트 다운로드 실패. 상태 코드: {response.status_code}")
    #else:
        #st.error(f"폰트 {font_name}이(가) 이미 존재합니다.")
        

    # 폰트 파일을 matplotlib에 적용
    font_prop = fm.FontProperties(fname=font_path)
    plt.rcParams['font.family'] = font_prop.get_name()
    #st.seccess(f"matplotlib에 폰트 {font_prop.get_name()} 적용 완료")

    return fm, plt

# 주어진 동영상 URL에서 자막을 추출하여 반환하는 함수.
def extract_text_from_video_url(video_url):
    try:
        loader = YoutubeLoader.from_youtube_url(video_url)
        transcript = loader.load()
        return transcript
    except Exception as e:
        st.error(f"동영상 자막을 추출하는 데 실패했습니다: {str(e)}")
        return None
        
# LLM을 통해 프롬프트와 파일을 전달하고 응답을 받는 함수
def run_llm_with_video_and_prompt(api_key, titles, requests, video_data_str):
    global global_generated_prompt
    openai.api_key = api_key

    responses = []
    global_generated_prompt = []  # 프롬프트들을 담을 리스트 초기화

    try:
        # 요청사항 리스트 문자열 생성
        request_list_str = "\n".join([
            f"{i+1}.{title}의 항목 데이터에 대해 '{request}' 요청 사항을 만족하게 답변해야 한다.\n"
            for i, (title, request) in enumerate(zip(titles, requests))
        ])

        # 프롬프트 텍스트 정의
        generated_prompt = f"""
        아애의 항목 데이터를 간결하고 깔끔한 보고서 작성을 위해 보고서 내용에 대해서 알기 쉽게 내용 요약하고 설명해야 한다.
        항목 데이터는 video의 자막 텍스트틀 추출한 데이터이니 이를 토대로 다음과 같은 조건에 모두 만족해야 한다.
        가. 아래의 항목 데이터를 분석하여 각 항목마다의 '요청사항' 리스트와 조건사항에 대해 모두 만족할 수 있도록 최적화된 보고서를 완성해.
        나. 항목 데이터 내 가장 첫번째 행은 각 보고서 항목에 타이틀이므로 순번과 문구를 그대로 유지해야 한다. 이 항목의 타이틀을 기준으로 각 항목 데이터를 분류하고 그에 맞는 요청사항을 반영해야 한다.
        다. 문단 끝날 때마다 줄바꿈을 해야 하고 줄바꿈은 <br/> 태그로 변환한다.
        라. 표 형식의로 table태그로 답변 할 때는 th과 td 태그는 border는 사이즈 1이고 색상은 검정색으로 구성한다. table 태그 가로길이는 전체를 차지해야 한다.
        마. 요약과 설명으로 답변 가시성 높게 특수기호를 활용하여 보고서 양식에 준하게 요약한 내용을 설명해줘야 하고, 보고서 양식에 맞춰 간결하고 깔끔하게 요약하고 html 형식으로 변환해야 한다.
        바. 이외 table 태그가 포함 안된 데이터는 파일에서 텍스트를 추출한 데이터로 내용으로 
        사. 답변할 때는 반드시 모든 항목 데이터의 수정한 데이터 내용과 HTML 형삭에 맞춰 답변한다. 문단마다 줄바꿈을 적용하여 br태그 활용하고 가시성 높게 특수기호를 활용하여 보고서 양식에 준하게 요약한 내용을 설명해줘야 한다.
        아. '✨AI 요약과 설명' 타이틀 추가하고 가장 먼저 나와야하고 위에는 그 어떤한 설명 내용도 먼저 응답하면 안 된다. 
        자. 전달받은 보고서 전반적인 내용에 대해 너가 선정한 가장 좋은 방법으로 요약과 설명하고 그 내용을 HTML 형식으로 변환하여 답변해야 한다.
        차. '🗣️번역 내용'이 하단에 추가하고 항목 데이터가 영어일 경우에는 요약과 설명을 한글로 번역하고 타이틀을 '🗣️번역 내용(한글)'로 출력하고
            한글일 경우에는 요역과 설명을 영어로 변역하고 타이틀은 '🗣️번역 내용(Endglish)'로 출력해야 한다.
        차. '````', '````HTML' 이 문구들이 답변에 포함되지 않아야 한다.
        -요청사항
        [
            {request_list_str}
        ]
        -항목 데이터
        [
            {video_data_str}
        ]
        """
        
        # 텍스트 길이 제한 확인 (예: 1000000자로 제한)
        if len(generated_prompt) > 1000000:
            st.error("프롬프트 글자 수 초과로 LLM 연동에 실패했습니다.")
        else:
            global_generated_prompt.append(generated_prompt)
            prompt_template = PromptTemplate(
                template=generated_prompt,
                input_variables=[]
            )

            # LLM 모델 생성
            llm = ChatOpenAI(model_name="gpt-4o")
            chain = LLMChain(llm=llm, prompt=prompt_template)

            success = False
            retry_count = 0
            max_retries = 5  # 최대 재시도 횟수

            # 응답을 받을 때까지 재시도
            while not success and retry_count < max_retries:
                try:
                    response = chain.run({})
                    responses.append(response)
                    success = True
                except RateLimitError:
                    retry_count += 1
                    st.warning(f"API 요청 한도를 초과했습니다. 10초 후 다시 시도합니다. 재시도 횟수: {retry_count}/{max_retries}")
                    time.sleep(10)

                time.sleep(10)
    except Exception as e:
        st.error(f"LLM 실행 중 오류가 발생했습니다: {str(e)}")

    return responses

# URL 패턴 검증을 위한 정규 표현식
def is_valid_url(url):
    url_pattern = re.compile(
        r'^(https?:\/\/)'  # http 또는 https로 시작
        r'((([A-Za-z]{1,})(\.[A-Za-z]{2,}))|([A-Za-z0-9.-]+\.[A-Za-z]{2,}))'  # 도메인 이름
        r'(:\d+)?(\/[A-Za-z0-9#_\/.-]*)?'  # 포트번호, 경로, 앵커, 기타 정보
        r'(\?[A-Za-z0-9=&]*)?$'  # 쿼리 파라미터
    )
    #URL이 유효한지 확인하는 함수
    return bool(url_pattern.match(url))

def extract_transcript_from_youtube(video_url):
    try:
        # YouTube 영상 ID 추출
        video_id = video_url.split('v=')[-1].split('&')[0]
        
        # YouTube 영상 자막 추출
        transcript = YouTubeTranscriptApi.get_transcript(video_id)
        return transcript
    except Exception as e:
        st.error(f"동영상 자막을 추출하는 데 실패했습니다: {str(e)}")
        return None
        
# Backend 기능 구현 끝 ---
