import os
import openai
import streamlit as st
from langchain.prompts import PromptTemplate
from langchain import LLMChain
import requests
import PyPDF2
import pandas as pd
import docx
import pptx
from PIL import Image
from io import BytesIO
import base64
import urllib.parse
from openai.error import RateLimitError
from langchain.chat_models import ChatOpenAI
import time
import json
import openpyxl
from openpyxl.utils import get_column_letter
import re
import speech_recognition as sr
from pydub import AudioSegment
import tempfile

# Backend 기능 구현 시작

# 전역변수로 프롬프트 및 파일 데이터 저장
global_generated_prompt = []
global_report_map = {}  # 전역 변수로 map 선언

# GitHub 정보 및 OpenAI API 키 자동 설정 또는 입력창을 통해 설정
def load_env_info():
    json_data = '''
    {
        "github_repo": "soryhon/hanachatbot",
        "github_branch": "main"
    }
    '''
    
    # JSON 데이터에서 정보 추출
    github_info = json.loads(json_data)
    github_repo = github_info['github_repo']
    github_branch = github_info['github_branch']
    
    # 환경 변수에서 GitHub Token 및 OpenAI API Key 가져오기
    github_token = os.getenv("GITHUB_TOKEN")
    openai_api_key = os.getenv("OPENAI_API_KEY")

    github_set = False
    openai_set = False

    # 입력창을 가로로 배치 (각각 50%의 너비)
    col1, col2 = st.columns(2)

    with col1:
        if not github_token:
            github_token = st.text_input("GitHub Token을 입력하세요", type="password", key="github_token_input")
        else:
            github_set = True
            st.success("GitHub 정보가 자동으로 설정되었습니다!")
        st.session_state["github_token"] = github_token

    with col2:
        if not openai_api_key:
            openai_api_key = st.text_input("OpenAI API 키를 입력하세요", type="password", key="openai_api_key_input")
        else:
            openai_set = True
            st.success("OpenAI API 키가 자동으로 설정되었습니다!")
        st.session_state["openai_api_key"] = openai_api_key

    # GitHub 저장소 정보 세션에 저장
    st.session_state["github_repo"] = github_repo
    st.session_state["github_branch"] = github_branch

    # GitHub 정보가 설정되었는지 확인하고 세션 상태 반영
    return github_set

# 다양한 파일 형식에서 데이터를 추출하는 함수
def extract_data_from_file(file_content, file_type):
    if file_content is None:
        st.error("파일 내용을 가져오지 못했습니다.")
        return None

    if file_type == 'pdf':
        return extract_text_from_pdf(file_content)
    elif file_type == 'csv':
        return extract_text_from_csv(file_content)
    elif file_type == 'docx':
        return extract_text_from_word(file_content)
    elif file_type == 'pptx':
        return extract_text_from_ppt(file_content)
    elif file_type in ['png', 'jpg', 'jpeg']:
        return extract_text_from_image(file_content)
    elif file_type == 'txt':
        return extract_text_from_txt(file_content)
    elif file_type == 'log':
        return extract_text_from_log(file_content)
    elif file_type == 'wav':  # 음성 파일 추가
        return extract_text_from_audio(file_content, file_type)    
    else:
        st.error(f"{file_type} 형식은 지원되지 않습니다.")
        return None

# PDF 파일에서 텍스트 추출
def extract_text_from_pdf(file_content):
    reader = PyPDF2.PdfReader(file_content)
    text = ''
    for page in range(len(reader.pages)):
        text += reader.pages[page].extract_text()
    return text

# CSV 파일에서 텍스트 추출
def extract_text_from_csv(file_content):
    csv_data = pd.read_csv(file_content)
    return csv_data

# 워드 파일에서 텍스트 추출
def extract_text_from_word(file_content):
    doc = docx.Document(file_content)
    return '\n'.join([para.text for para in doc.paragraphs])

# PPT 파일에서 텍스트 추출
def extract_text_from_ppt(file_content):
    presentation = pptx.Presentation(file_content)
    text = ''
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + '\n'
    return text

# 텍스트 파일에서 텍스트 추출
def extract_text_from_txt(file_content):
    try:
        # BytesIO 객체의 내용을 문자열로 변환
        if isinstance(file_content, BytesIO):
            return file_content.read().decode("utf-8")
        elif isinstance(file_content, str):
            return file_content
        else:
            st.error("알 수 없는 파일 형식입니다.")
            return None
    except Exception as e:
        st.error(f"txt 파일에서 텍스트를 추출하는 중 오류가 발생했습니다: {str(e)}")
        return None

# 로그 파일에서 텍스트 추출
def extract_text_from_log(file_content):
    try:
        # BytesIO 객체의 내용을 문자열로 변환
        if isinstance(file_content, BytesIO):
            return file_content.read().decode("utf-8")
        elif isinstance(file_content, str):
            return file_content
        else:
            st.error("알 수 없는 파일 형식입니다.")
            return None
    except Exception as e:
        st.error(f"log 파일에서 텍스트를 추출하는 중 오류가 발생했습니다: {str(e)}")
        return None

    
# 이미지에서 텍스트 추출 (OCR)
def extract_text_from_image(file_content):
    image = Image.open(file_content)
    return "이미지에서 텍스트를 추출하는 기능은 구현되지 않았습니다."


# GitHub에 폴더가 존재하는지 확인하고 없으면 생성하는 함수
def create_github_folder_if_not_exists(repo, folder_name, token, branch='main'):
    url = f"https://api.github.com/repos/{repo}/contents/{folder_name}?ref={branch}"
    headers = {"Authorization": f"token {token}"}
    response = requests.get(url, headers=headers)
    
    if response.status_code == 404:
        st.warning(f"'{folder_name}' 폴더가 존재하지 않아 생성 중입니다.")
        create_folder_url = f"https://api.github.com/repos/{repo}/contents/{folder_name}"
        data = {
            "message": f"Create {folder_name} folder",
            "content": base64.b64encode(b'').decode('utf-8'),  # 빈 파일로 폴더 생성
            "branch": branch
        }
        requests.put(create_folder_url, json=data, headers=headers)
        st.success(f"'{folder_name}' 폴더가 성공적으로 생성되었습니다.")
    # 폴더가 이미 존재할 경우 메시지를 표시하지 않음

# GitHub API 요청을 처리하는 함수 (파일 목록을 가져옴)
def get_github_files(repo, branch, token, folder_name='uploadFiles'):
    create_github_folder_if_not_exists(repo, folder_name, token, branch)
    url = f"https://api.github.com/repos/{repo}/git/trees/{branch}?recursive=1"
    headers = {"Authorization": f"token {token}"}
    response = requests.get(url, headers=headers)
    if response.status_code == 200:
        files = [item['path'] for item in response.json().get('tree', []) if item['type'] == 'blob' and item['path'].startswith(folder_name)]
        return files
    else:
        st.error("GitHub 파일 목록을 가져오지 못했습니다. 저장소 정보나 토큰을 확인하세요.")
        return []

# GitHub에서 파일의 SHA 값을 가져오는 함수
def get_file_sha(repo, file_path, token, branch='main'):
    encoded_file_path = urllib.parse.quote(file_path)
    url = f"https://api.github.com/repos/{repo}/contents/{encoded_file_path}?ref={branch}"
    headers = {"Authorization": f"token {token}"}
    response = requests.get(url, headers=headers)
    
    if response.status_code == 200:
        return response.json().get('sha', None)
    else:
        return None

# GitHub에 파일 업로드 함수 (덮어쓰기 포함)
def upload_file_to_github(repo, folder_name, file_name, file_content, token, branch='main', sha=None):
    create_github_folder_if_not_exists(repo, folder_name, token, branch)
    encoded_file_name = urllib.parse.quote(file_name)
    url = f"https://api.github.com/repos/{repo}/contents/{folder_name}/{encoded_file_name}"
    headers = {
        "Authorization": f"token {token}",
        "Content-Type": "application/json"
    }

    content_encoded = base64.b64encode(file_content).decode('utf-8')

    data = {
        "message": f"Upload {file_name}",
        "content": content_encoded,
        "branch": branch
    }

    if sha:
        data["sha"] = sha

    response = requests.put(url, json=data, headers=headers)

    if response.status_code in [200, 201]:
        st.success(f"{file_name} 파일이 성공적으로 업로드(덮어쓰기) 되었습니다.")
    else:
        st.error(f"파일 업로드에 실패했습니다: {response.status_code}")
        st.error(response.json())

# GitHub에서 파일을 다운로드하는 함수
def get_file_from_github(repo, branch, filepath, token):
    encoded_filepath = urllib.parse.quote(filepath)
    url = f"https://api.github.com/repos/{repo}/contents/{encoded_filepath}?ref={branch}"
    headers = {"Authorization": f"token {token}"}
    response = requests.get(url, headers=headers)
    
    if response.status_code == 200:
        return BytesIO(requests.get(response.json()['download_url']).content)
    else:
        st.error(f"{filepath} 파일을 가져오지 못했습니다. 상태 코드: {response.status_code}")
        return None

# 엑셀 시트에서 셀 스타일 및 병합 정보 추출
def extract_cell_style_and_merged(ws):
    style_dict = {}
    merged_cells = ws.merged_cells.ranges
    
    for row in ws.iter_rows():
        for cell in row:
            cell_style = {
                "alignment": cell.alignment.horizontal if cell.alignment else 'left',
                "font_bold": cell.font.bold if cell.font else False,
                "border": bool(cell.border and (cell.border.left or cell.border.right or cell.border.top or cell.border.bottom)),
                "value": "" if str(cell.value).startswith("Unnamed") else cell.value
            }
            style_dict[cell.coordinate] = cell_style
            
    return style_dict, merged_cells

# 엑셀 시트 데이터를 HTML로 변환하고 스타일 및 병합 적용
def convert_df_to_html_with_styles_and_merging(ws, df):
    # 1행의 'Unnamed: 숫자' 형식 값은 공백으로 처리
    df.columns = [re.sub(r'Unnamed: \d+', '', str(col)).strip() for col in df.columns]
    
    style_dict, merged_cells = extract_cell_style_and_merged(ws)
    df = df.fillna('')  # NaN 값을 공백으로 처리
    html = "<table class='table table-bordered' style='border-spacing: 0; border-collapse: collapse;'>\n"

    # 헤더 부분 (border 간격 없게 설정)
    html += "<thead>\n<tr>\n"
    for col in df.columns:
        html += f"<th style='text-align:center; font-weight:bold; background-color:#E7E6E6; border: 1px solid black;'>{col}</th>\n"
    html += "</tr>\n</thead>\n"

    # 병합 셀 정보 저장
    merged_cells_dict = {}
    for merged in merged_cells:
        min_row, min_col, max_row, max_col = merged.min_row, merged.min_col, merged.max_row, merged.max_col
        for row in range(min_row, max_row + 1):
            for col in range(min_col, max_col + 1):
                merged_cells_dict[(row, col)] = (min_row, min_col, max_row, max_col)

    # 데이터 부분
    html += "<tbody>\n"
    for row_idx, row in df.iterrows():
        html += "<tr>\n"
        for col_idx, value in enumerate(row):
            cell_ref = f"{get_column_letter(col_idx+1)}{row_idx+2}"  # 셀 참조 계산
            style = style_dict.get(cell_ref, {})
            alignment = style.get("alignment", "left")
            font_weight = "bold" if style.get("font_bold", False) else "normal"
            border = "1px solid black" if style.get("border", False) else "none"
            cell_coordinates = (row_idx + 2, col_idx + 1)

            if cell_coordinates in merged_cells_dict:
                min_row, min_col, max_row, max_col = merged_cells_dict[cell_coordinates]
                rowspan = max_row - min_row + 1
                colspan = max_col - min_col + 1
                if (row_idx + 2, col_idx + 1) == (min_row, min_col):
                    html += f"<td rowspan='{rowspan}' colspan='{colspan}' style='text-align:{alignment}; font-weight:{font_weight}; border:{border};'>{value}</td>\n"
            elif cell_coordinates not in merged_cells_dict:
                html += f"<td style='text-align:{alignment}; font-weight:{font_weight}; border:{border};'>{value}</td>\n"
        html += "</tr>\n"
    html += "</tbody>\n</table>"
    
    return html

# 엑셀 파일에서 시트 가져오기 및 데이터 추출
def extract_sheets_from_excel(file_content, selected_sheets):
    try:
        excel_data = pd.ExcelFile(file_content)
        all_sheets = excel_data.sheet_names
        
        # 사용자가 선택한 시트를 가져옴
        if selected_sheets == 'all':
            selected_sheets = all_sheets
        else:
            selected_sheets = [all_sheets[int(i)-1] for i in selected_sheets if int(i) <= len(all_sheets)]
        
        # 선택한 시트의 데이터를 하나로 합침
        data_dict = {}
        with pd.ExcelFile(file_content) as xls:
            for sheet in selected_sheets:
                df = pd.read_excel(xls, sheet_name=sheet)
                data_dict[sheet] = df
                
        return data_dict
    except Exception as e:
        st.error(f"엑셀 파일의 시트 데이터를 추출하는 중에 오류가 발생했습니다: {str(e)}")
        return None

# 숫자, ',' 및 '-'만 허용하는 함수
def validate_sheet_input(input_value):
    if all(c.isdigit() or c in ['-', ','] for c in input_value):
        return True
    return False

# 시트 선택 로직 추가
def handle_sheet_selection(file_content, sheet_count, idx):
    # 3개의 객체를 가로로 배치
    col1, col2, col3 = st.columns([0.33, 0.33, 0.33])
    
    with col1:
        st.text_input(f"시트 갯수_{idx}", value=f"{sheet_count}개", disabled=True)  # 시트 갯수 표시 (비활성화)
    
    with col2:
        sheet_selection = st.text_input(f"시트 선택_{idx}(예: 1-3, 5)", value="1", key=f"sheet_selection_{idx}")

    with col3:
        select_button = st.button("선택", key=f"select_button_{idx}")

    # 시트 선택 버튼이 눌렸을 때만 파일 데이터를 가져옴
    if select_button:
        if validate_sheet_input(sheet_selection):
            selected_sheets = parse_sheet_selection(sheet_selection, sheet_count)
            if selected_sheets:
                file_data = extract_sheets_from_excel(file_content, selected_sheets)
                return file_data
            else:
                st.error("선택한 시트가 잘못되었습니다.")
        else:
            st.error("잘못된 입력입니다. 숫자와 '-', ',' 만 입력할 수 있습니다.")
    return None

# 시트 선택 입력값을 분석하는 함수
def parse_sheet_selection(selection, sheet_count):
    selected_sheets = []

    try:
        if '-' in selection:
            start, end = map(int, selection.split('-'))
            if start <= end <= sheet_count:
                selected_sheets.extend(list(range(start, end+1)))
        elif ',' in selection:
            selected_sheets = [int(i) for i in selection.split(',') if 1 <= int(i) <= sheet_count]
        else:
            selected_sheets = [int(selection)] if 1 <= int(selection) <= sheet_count else []
    except ValueError:
        st.error("잘못된 시트 선택 입력입니다.")
        return None

    return selected_sheets

# 파일에서 데이터를 추출하고 요청사항 리스트에서 선택한 엑셀 파일의 시트를 보여주는 로직 수정
def handle_file_selection(file_path, file_content, file_type, idx):
    if file_type == 'xlsx':
        wb = openpyxl.load_workbook(file_content)
        sheet_count = len(wb.sheetnames)

        # 시트 선택 로직 처리
        file_data_dict = handle_sheet_selection(file_content, sheet_count, idx)
        return file_data_dict
    else:
        return extract_data_from_file(file_content, file_type)

# HTML 보고서 생성 함수 (배열에서 데이터 가져옴)
def generate_final_html_report():
    report_html = ""
    if len(global_report_map) > 0:  # map 변수가 null이 아니고 사이즈가 1 이상일 때
        for idx, file_data in global_report_map.items():  # map의 데이터를 가져옴
            if file_data:
                report_html += f"<div style='text-indent: 1px;'>\n{file_data}\n</div>\n"
                #report_html += f"{file_data}\n"
                report_html += f"{idx}--<p/>"  # 줄바꿈 추가

        st.session_state['html_report'] = report_html  # 최종 값을 세션 상태에 저장

# 엑셀 데이터 및 제목을 HTML로 변환하여 하나의 세트로 출력하는 함수
def generate_html_report_with_title(titles, data_dicts):
    report_html = ""
    
    for i, (title, data_dict) in enumerate(zip(titles, data_dicts), start=1):
        report_html += f"<h3>{i}. {title}</h3>\n"
        report_html += "<div style='text-indent: 20px;'>\n"
        
        for sheet_name, df in data_dict.items():
            wb = openpyxl.load_workbook(BytesIO(df))
            ws = wb[sheet_name]
            report_html += convert_df_to_html_with_styles_and_merging(ws, df)
        
        report_html += "</div>\n"
    
    return report_html

# LLM을 통해 프롬프트와 파일을 전달하고 응답을 받는 함수
def run_llm_with_file_and_prompt(api_key, titles, requests, file_data_list):
    global global_generated_prompt
    openai.api_key = api_key

    responses = []
    global_generated_prompt = []  # 프롬프트들을 담을 리스트 초기화

    try:
        for i, (title, request, file_data) in enumerate(zip(titles, requests, file_data_list)):
            if isinstance(file_data, pd.DataFrame):
                file_data_str = file_data.to_string()
            elif isinstance(file_data, dict):
                file_data_str = "\n".join([f"시트 이름: {sheet_name}\n{data.to_string()}" for sheet_name, data in file_data.items()])
            else:
                file_data_str = str(file_data)

            # 프롬프트 생성
            generated_prompt = f"""
            보고서 제목은 '{title}'로 하고, 아래의 파일 데이터를 분석하여 '{request}'를 요구 사항을 만족할 수 있도록 최적화된 보고서를 완성해.
            표로 표현 할 때는 table 태그 형식으로 구현해야 한다. th과 td 태그는 border는 사이즈 1이고 색상은 검정색으로 구성한다.
            표의 첫번째 행은 타이틀이 이므로 th태그로 구현하고 가운데 정렬, bold처리 해야 한다. 바탕색은 '#E7E6E6' 이어야 한다.
            예시와 같은 구조로 구성한다. 보고서 제목은 앞에 순번을 표시하고 바로 아래 요구 사항에 맞는 내용을 이어서 보여줘야 한다.
            예시 : '\r\n {i+1}. 보고서 제목\r\n보고서 내용'
            파일 데이터: {file_data_str}
            """

            global_generated_prompt.append(generated_prompt)

            prompt_template = PromptTemplate(
                template=generated_prompt,
                input_variables=[]
            )

            # LLM 모델 생성
            llm = ChatOpenAI(model_name="gpt-4")
            chain = LLMChain(llm=llm, prompt=prompt_template)

            success = False
            retry_count = 0
            max_retries = 5  # 최대 재시도 횟수

            # 응답을 받을 때까지 재시도
            while not success and retry_count < max_retries:
                try:
                    response = chain.run({})
                    responses.append(response)
                    success = True
                except RateLimitError:
                    retry_count += 1
                    st.warning(f"API 요청 한도를 초과했습니다. 10초 후 다시 시도합니다. 재시도 횟수: {retry_count}/{max_retries}")
                    time.sleep(10)

                time.sleep(10)
    except Exception as e:
        st.error(f"LLM 실행 중 오류가 발생했습니다: {str(e)}")
    return responses

# Backend 기능 구현 끝 

# Front 기능 구현 시작

# GitHub 정보가 있는지 확인하고 파일 업로드 객체를 출력
github_info_loaded = load_env_info()

# 업로드 가능한 파일 크기 제한 (100MB)
MAX_FILE_SIZE_MB = 100
MAX_FILE_SIZE_BYTES = 100 * 1024 * 1024

# 1 프레임
# 파일 업로드
st.subheader("1. 파일 업로드")

# 지원되는 파일 형식 리스트
supported_file_types = ['xlsx', 'pptx', 'docx', 'csv', 'png', 'jpg', 'jpeg', 'pdf', 'txt', 'log']

if github_info_loaded:
    with st.expander("파일 업로드", expanded=True):
        uploaded_files = st.file_uploader("파일을 여러 개 드래그 앤 드롭하여 업로드하세요. (최대 100MB)", accept_multiple_files=True)

        if uploaded_files:
            for uploaded_file in uploaded_files:
                file_type = uploaded_file.name.split('.')[-1].lower()

                if file_type not in supported_file_types:
                    st.error(f"지원하지 않는 파일입니다: {uploaded_file.name}")
                    continue

                if uploaded_file.size > MAX_FILE_SIZE_BYTES:
                    st.warning(f"'{uploaded_file.name}' 파일은 {MAX_FILE_SIZE_MB}MB 제한을 초과했습니다. 파일 크기를 줄이거나 GitHub에 직접 푸시하세요.")
                else:
                    file_content = uploaded_file.read()
                    file_name = uploaded_file.name
                    folder_name = 'uploadFiles'

                    sha = get_file_sha(st.session_state['github_repo'], f"{folder_name}/{file_name}", st.session_state['github_token'], branch=st.session_state['github_branch'])

                    if sha:
                        st.warning(f"'{file_name}' 파일이 이미 존재합니다. 덮어쓰시겠습니까?")
                        col1, col2 = st.columns(2)

                        with col1:
                            if st.button(f"'{file_name}' 덮어쓰기", key=f"overwrite_{file_name}"):
                                upload_file_to_github(st.session_state['github_repo'], folder_name, file_name, file_content, st.session_state['github_token'], branch=st.session_state['github_branch'], sha=sha)
                                st.success(f"'{file_name}' 파일이 성공적으로 덮어쓰기 되었습니다.")
                                uploaded_files = None
                                break

                        with col2:
                            if st.button("취소", key=f"cancel_{file_name}"):
                                st.info("덮어쓰기가 취소되었습니다.")
                                uploaded_files = None
                                break
                    else:
                        upload_file_to_github(st.session_state['github_repo'], folder_name, file_name, file_content, st.session_state['github_token'])
                        st.success(f"'{file_name}' 파일이 성공적으로 업로드되었습니다.")
                        uploaded_files = None
else:
    st.warning("GitHub 정보가 저장되기 전에는 파일 업로드를 할 수 없습니다. 먼저 GitHub 정보를 입력해 주세요.")

# 2 프레임: 작성 보고서 요청사항
st.subheader("2. 작성 보고서 요청사항")

# 요청사항 리스트
with st.expander("요청사항 리스트", expanded=True):
    if 'rows' not in st.session_state:
        st.session_state['rows'] = [{"제목": "", "요청": "", "파일": "", "데이터": "", "checked": False,"파일데이터":""}]

    rows = st.session_state['rows']
    checked_rows = []

    for idx, row in enumerate(rows):
        with st.container():
            col1, col2 = st.columns([0.05, 0.95])  # 체크박스와 제목 부분을 가로로 나눔
            with col1:
                row_checked = st.checkbox("", key=f"row_checked_{idx}", value=row.get("checked", False))  # 체크박스만 추가
            with col2:
                st.markdown(f"#### 요청사항 {idx+1}")
        
            row['제목'] = st.text_input(f"제목_{idx} (요청사항 {idx+1})", row['제목'], key=f"title_{idx}")
            row['요청'] = st.text_area(f"요청_{idx} (요청사항 {idx+1})", row['요청'], key=f"request_{idx}")

            file_list = ['파일을 선택하세요.']
            if st.session_state.get('github_token') and st.session_state.get('github_repo'):
                file_list += get_github_files(st.session_state['github_repo'], st.session_state['github_branch'], st.session_state['github_token'])

            selected_file = st.selectbox(f"파일 선택_{idx} (요청사항 {idx+1})", options=file_list, key=f"file_select_{idx}")

            if selected_file != '파일을 선택하세요.':
                file_path = selected_file
                file_content = get_file_from_github(st.session_state["github_repo"], st.session_state["github_branch"], file_path, st.session_state["github_token"])
                
                if file_content:
                    file_type = file_path.split('.')[-1].lower()

                    # 파일 형식 검증 (지원되는 파일만 처리)
                    if file_type not in supported_file_types:
                        st.error(f"지원하지 않는 파일입니다: {file_path}")
                        row['데이터'] = ""
                    else:
                        # 엑셀 파일인 경우 시트 선택 로직을 추가
                        if file_type == 'xlsx':
                            html_report_set = f"<div style='text-indent: 5px;'>\n"
                            file_data_dict = handle_file_selection(file_path, file_content, file_type, idx)
                            if file_data_dict is not None:
                                
                                # 제목 입력 값 가져오기
                                html_report_set +=  f"<h3>{idx + 1}. {row['제목']}</h3>\n"
                                row['파일'] = f"/{st.session_state['github_repo']}/{st.session_state['github_branch']}/{selected_file}"
                                for sheet_name, df in file_data_dict.items():
                                    wb = openpyxl.load_workbook(file_content)
                                    ws = wb[sheet_name]
                                    html_data = convert_df_to_html_with_styles_and_merging(ws, df)
                                    html_report_set += f"<div style='text-indent: 20px;'>{html_data}</div>\n"
                        else:                           
                            file_data = extract_data_from_file(file_content, file_type)
                            if file_data:     
                                html_report_set = f"<div style='text-indent: 5px;'>\n"
                                # 제목 입력 값 가져오기
                                html_report_set +=  f"<h3>{idx + 1}. {row['제목']}</h3>\n"
                                html_report_set += f"<p>{file_data}</p>"                        
                        html_report_set += "</div>\n"       
                        row['파일데이터'] = html_report_set
                        global_report_map[idx] = html_report_set

                        #if "html_report" in st.session_state:
                        generate_final_html_report()

                else:
                    st.error(f"{selected_file} 파일을 GitHub에서 불러오지 못했습니다.")
                
            st.text_input(f"파일 경로_{idx} (요청사항 {idx+1})", row['파일'], disabled=True, key=f"file_{idx}")

        if row_checked:
            checked_rows.append(idx)

    # 행 추가 및 삭제 버튼을 가로로 배치하고 각 버튼의 너비를 30%로 설정
    col1, col2, col3 = st.columns([0.3, 0.3, 0.3])

    with col1:
        if st.button("행 추가", key="add_row", use_container_width=True):
            new_row = {"제목": "", "요청": "", "파일": "", "데이터": "", "checked": False,"파일데이터":""}
            st.session_state['rows'].append(new_row)

    with col2:
        if st.button("행 삭제", key="delete_row", use_container_width=True):
            if checked_rows:
                st.session_state['rows'] = [row for idx, row in enumerate(rows) if idx not in checked_rows]
                st.success(f"체크된 {len(checked_rows)}개의 요청사항이 삭제되었습니다.")
            else:
                st.warning("삭제할 요청사항을 선택해주세요.")
    
    with col3:
        if st.button("새로고침", key="refresh_page", use_container_width=True):
            st.success("새로고침 하였습니다.")

# 3 프레임
# 보고서 작성 실행 버튼
st.subheader("3. 보고서 작성 실행")

if st.button("보고서 작성", key="generate_report", use_container_width=True):
    if not st.session_state.get("openai_api_key"):
        st.error("먼저 OpenAI API 키를 입력하고 저장하세요!")
    elif not st.session_state['rows'] or all(not row["제목"] or not row["요청"] or not row["데이터"] for row in st.session_state['rows']):
        st.error("요청사항의 제목, 요청, 파일을 모두 입력해야 합니다!")
    else:
        titles = [row['제목'] for row in st.session_state['rows']]
        requests = [row['요청'] for row in st.session_state['rows']]
        file_data_list = [row['데이터'] for row in st.session_state['rows']]

        responses = run_llm_with_file_and_prompt(
            st.session_state["openai_api_key"], 
            titles, 
            requests, 
            file_data_list
        )
        st.session_state["response"] = responses

# 양식 저장, 양식 불러오기 버튼을 같은 행에 가로로 배치하고 각 버튼의 너비를 50%로 설정
col1, col2 = st.columns([0.5, 0.5])
with col1:
    if st.button("양식 저장", key="save_template", use_container_width=True):
        st.success("양식이 저장되었습니다.")

with col2:
    if st.button("양식 불러오기", key="load_template", use_container_width=True):
        st.success("양식이 불러와졌습니다.")


# 4 프레임: 결과 보고서
st.subheader("4. 결과 보고서")

# 결과 보고서 HTML 보기
if "html_report" in st.session_state:
   #generate_final_html_report()
    st.components.v1.html(st.session_state['html_report'], height=1280, scrolling=True)

# 전달된 프롬프트
st.text_area("전달된 프롬프트:", value="\n\n".join(global_generated_prompt), height=150)

# LLM 응답 보기
if "response" in st.session_state:
    for idx, response in enumerate(st.session_state["response"]):
        st.text_area(f"응답 {idx+1}:", value=response, height=300)
        st.components.v1.html(response, height=600, scrolling=True)

# Front 기능 구현 끝
