import os
import openai
import streamlit as st
from langchain.prompts import PromptTemplate
from langchain import LLMChain
import requests
import PyPDF2
import pandas as pd
import docx
import pptx
from PIL import Image
from io import BytesIO

# 페이지 너비를 전체 화면으로 설정
st.set_page_config(layout="wide")

# GitHub API 요청을 처리하는 함수
def get_github_files(repo, branch, token):
    url = f"https://api.github.com/repos/{repo}/git/trees/{branch}?recursive=1"
    headers = {"Authorization": f"token {token}"}
    response = requests.get(url, headers=headers)
    if response.status_code == 200:
        files = [item['path'] for item in response.json().get('tree', []) if item['type'] == 'blob']
        return files
    else:
        st.error("GitHub 파일 목록을 가져오지 못했습니다. 저장소 정보나 토큰을 확인하세요.")
        return []

# GitHub에서 파일을 다운로드하는 함수
def get_file_from_github(repo, branch, filepath, token):
    url = f"https://api.github.com/repos/{repo}/contents/{filepath}?ref={branch}"
    headers = {"Authorization": f"token {token}"}
    response = requests.get(url, headers=headers)
    
    if response.status_code == 200:
        return BytesIO(requests.get(response.json()['download_url']).content)
    else:
        st.error(f"{filepath} 파일을 가져오지 못했습니다.")
        return None

# 다양한 파일 형식에서 데이터를 추출하는 함수
def extract_data_from_file(file_content, file_type):
    if file_type == 'pdf':
        return extract_text_from_pdf(file_content)
    elif file_type == 'xlsx':
        return extract_text_from_excel(file_content)
    elif file_type == 'docx':
        return extract_text_from_word(file_content)
    elif file_type == 'pptx':
        return extract_text_from_ppt(file_content)
    elif file_type in ['png', 'jpg', 'jpeg']:
        return extract_text_from_image(file_content)
    else:
        st.error(f"{file_type} 형식은 지원되지 않습니다.")
        return None

# PDF 파일에서 텍스트 추출
def extract_text_from_pdf(file_content):
    reader = PyPDF2.PdfReader(file_content)
    text = ''
    for page in range(len(reader.pages)):
        text += reader.pages[page].extract_text()
    return text

# 엑셀 파일에서 텍스트 추출
def extract_text_from_excel(file_content):
    excel_data = pd.read_excel(file_content)
    return excel_data.to_string()

# 워드 파일에서 텍스트 추출
def extract_text_from_word(file_content):
    doc = docx.Document(file_content)
    return '\n'.join([para.text for para in doc.paragraphs])

# PPT 파일에서 텍스트 추출
def extract_text_from_ppt(file_content):
    presentation = pptx.Presentation(file_content)
    text = ''
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + '\n'
    return text

# 이미지에서 텍스트 추출 (OCR)
def extract_text_from_image(file_content):
    image = Image.open(file_content)
    return "이미지에서 텍스트를 추출하는 기능은 구현되지 않았습니다."

# 프롬프트를 구성하고 LLM에 전달하는 함수
def create_prompt_and_send_to_llm(api_key, title, request, file_data):
    openai.api_key = api_key

    # 프롬프트 템플릿 구성
    prompt = f"""
    보고서 제목은 '{title}'로 하고, 이 파일에서 '{request}'를 만족할 수 있도록 최적화된 보고서를 완성해.
    파일 데이터: {file_data}
    """

    # LLM에 프롬프트 전달하고 응답 받기
    response = openai.ChatCompletion.create(
        model="gpt-3.5-turbo",
        messages=[
            {"role": "system", "content": "You are a helpful assistant."},
            {"role": "user", "content": prompt}
        ],
        max_tokens=1500
    )

    return response['choices'][0]['message']['content']

# 1. 프레임: GitHub 정보 저장 및 OpenAI API 키 저장
st.subheader("1. GitHub 정보 저장 및 OpenAI API 키 저장")

col1, col2 = st.columns(2)  # 두 개의 컬럼으로 나눔

# GitHub 정보 저장
with col1:
    st.subheader("GitHub 정보 입력")
    repo_input = st.text_input("GitHub 저장소 (owner/repo 형식)", value="")  # 공백으로 설정
    branch_input = st.text_input("GitHub 브랜치", value="main")
    token_input = st.text_input("GitHub Token", type="password")

    if st.button("GitHub 정보 저장"):
        if repo_input and branch_input and token_input:
            st.session_state["github_repo"] = repo_input
            st.session_state["github_branch"] = branch_input
            st.session_state["github_token"] = token_input
            st.success("GitHub 정보가 성공적으로 저장되었습니다!")
        else:
            st.error("모든 GitHub 정보를 입력해야 합니다!")

# OpenAI API 키 저장
with col2:
    st.subheader("OpenAI API 키 저장")
    api_key_input = st.text_input("OpenAI API 키를 입력하세요", type="password")
    if st.button("API 키 저장"):
        if api_key_input:
            st.session_state["api_key"] = api_key_input
            os.environ["OPENAI_API_KEY"] = api_key_input
            st.success("API 키가 성공적으로 저장되었습니다!")
        else:
            st.error("API 키를 입력해야 합니다!")

# 2. 프레임: 작성 보고서 요청사항과 실행 버튼
st.subheader("2. 작성 보고서 요청사항 및 실행 버튼")

col1, col2 = st.columns([0.8, 0.2])  # 가로 길이 80%, 20%

# 작성 보고서 요청사항 테이블
with col1:
    st.write("제목")
    title = st.text_input("", value="", disabled=False)

    st.write("요청")
    request = st.text_area("", disabled=False)

    st.write("파일")
    selected_file = None
    if "github_token" in st.session_state:
        files = get_github_files(st.session_state["github_repo"], st.session_state["github_branch"], st.session_state["github_token"])
        if files:
            selected_file = st.selectbox("GitHub 파일을 선택하세요", ["파일을 선택하세요"] + files, index=0)
        else:
            st.info("저장소에 파일이 없습니다.")
    
    st.write("데이터")
    if selected_file and selected_file != "파일을 선택하세요":
        file_path = selected_file
        file_content = get_file_from_github(st.session_state["github_repo"], st.session_state["github_branch"], file_path, st.session_state["github_token"])
        file_type = file_path.split('.')[-1].lower()
        file_data = extract_data_from_file(file_content, file_type)
        st.text_input("", value=f"선택한 파일 경로: {selected_file}", disabled=True)

# 실행 버튼
with col2:
    st.write(" ")
    st.write(" ")
    if st.button("실행"):
        if not st.session_state.get("api_key"):
            st.error("먼저 OpenAI API 키를 입력하고 저장하세요!")
        elif not title or not request or selected_file == "파일을 선택하세요":
            st.error("제목, 요청사항, 또는 파일을 선택해야 합니다!")
        else:
            response = create_prompt_and_send_to_llm(st.session_state["api_key"], title, request, file_data)
            st.session_state["response"] = response

# 3. 프레임: 결과 보고서
st.subheader("3. 결과 보고서")
if "response" in st.session_state:
    st.text_area("응답:", value=st.session_state["response"], height=300)
