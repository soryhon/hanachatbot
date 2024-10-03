import os
import openai
import streamlit as st
from langchain.prompts import PromptTemplate
from langchain import LLMChain
import requests
import PyPDF2
import pandas as pd
import docx
import pptx
from PIL import Image
from io import BytesIO
import base64
import time
from openai.error import RateLimitError
from langchain.chat_models import ChatOpenAI

# 전역변수로 프롬프트 저장
global_generated_prompt = []

# 페이지 너비를 전체 화면으로 설정
st.set_page_config(layout="wide")

# GitHub API 요청을 처리하는 함수
def get_github_files(repo, branch, token, folder_name=None):
    url = f"https://api.github.com/repos/{repo}/contents/{folder_name}?ref={branch}"
    headers = {"Authorization": f"token {token}"}
    response = requests.get(url, headers=headers)
    if response.status_code == 200:
        files = [item['name'] for item in response.json() if item['type'] == 'file']
        return files
    else:
        return []

# GitHub에서 파일의 SHA 값을 가져오는 함수
def get_file_sha(repo, file_path, token, branch='main'):
    url = f"https://api.github.com/repos/{repo}/contents/{file_path}?ref={branch}"
    headers = {"Authorization": f"token {token}"}
    response = requests.get(url, headers=headers)
    
    if response.status_code == 200:
        return response.json().get('sha', None)
    else:
        return None

# GitHub에 파일 업로드 함수
def upload_file_to_github(repo, folder_name, file_name, file_content, token, branch='main', sha=None):
    url = f"https://api.github.com/repos/{repo}/contents/{folder_name}/{file_name}"
    headers = {
        "Authorization": f"token {token}",
        "Content-Type": "application/json"
    }

    content_encoded = base64.b64encode(file_content).decode('utf-8')

    data = {
        "message": f"Upload {file_name}",
        "content": content_encoded,
        "branch": branch
    }

    if sha:
        data["sha"] = sha

    response = requests.put(url, json=data, headers=headers)

    if response.status_code == 201:
        st.success(f"{file_name} 파일이 성공적으로 업로드되었습니다.")
    elif response.status_code == 200:
        st.success(f"{file_name} 파일이 성공적으로 덮어쓰기 되었습니다.")
    else:
        st.error(f"파일 업로드에 실패했습니다: {response.status_code}")
        st.error(response.json())

# GitHub에서 파일을 다운로드하는 함수
def get_file_from_github(repo, branch, filepath, token):
    url = f"https://api.github.com/repos/{repo}/contents/{filepath}?ref={branch}"
    headers = {"Authorization": f"token {token}"}
    response = requests.get(url, headers=headers)

    if response.status_code == 404:
        st.error(f"{filepath} 파일을 가져오지 못했습니다. 상태 코드: 404 (파일을 찾을 수 없습니다). 경로를 확인하세요.")
        return None
    elif response.status_code != 200:
        st.error(f"{filepath} 파일을 가져오지 못했습니다. 상태 코드: {response.status_code}")
        return None

    try:
        return BytesIO(requests.get(response.json()['download_url']).content)
    except Exception as e:
        st.error(f"파일 데이터를 처리하는 중 오류가 발생했습니다: {str(e)}")
        return None

# 다양한 파일 형식에서 데이터를 추출하는 함수
def extract_data_from_file(file_content, file_type):
    if file_type == 'pdf':
        return extract_text_from_pdf(file_content)
    elif file_type == 'xlsx':
        return extract_text_from_excel(file_content)
    elif file_type == 'csv':
        return extract_text_from_csv(file_content)
    elif file_type == 'docx':
        return extract_text_from_word(file_content)
    elif file_type == 'pptx':
        return extract_text_from_ppt(file_content)
    elif file_type in ['png', 'jpg', 'jpeg']:
        return extract_text_from_image(file_content)
    else:
        st.error(f"{file_type} 형식은 지원되지 않습니다.")
        return None

# PDF 파일에서 텍스트 추출
def extract_text_from_pdf(file_content):
    try:
        if file_content is None:
            raise ValueError("유효하지 않은 파일입니다. PDF 파일이 없습니다.")
        
        file_content.seek(0)
        reader = PyPDF2.PdfReader(file_content)
        text = ''
        for page in range(len(reader.pages)):
            text += reader.pages[page].extract_text()
        return text
    except Exception as e:
        st.error(f"PDF 파일을 처리하는 중 오류가 발생했습니다: {str(e)}")
        return None

# 엑셀 파일에서 텍스트 추출 (시트 정보를 정확하게 가져오도록 수정)
def extract_text_from_excel(file_content):
    try:
        if file_content is None:
            raise ValueError("유효하지 않은 파일입니다. 파일이 없거나 경로가 잘못되었을 수 있습니다.")
        excel_data = pd.read_excel(file_content, sheet_name=None)  # 모든 시트를 불러옴
        all_data = {}

        for sheet_name, data in excel_data.items():
            all_data[sheet_name] = data

        return all_data  # 각 시트 이름과 데이터가 포함된 딕셔너리 반환
    except Exception as e:
        st.error(f"엑셀 파일을 불러오는 중 오류가 발생했습니다: {str(e)}")
        return None

# CSV 파일에서 텍스트 추출
def extract_text_from_csv(file_content):
    csv_data = pd.read_csv(file_content)
    return csv_data

# 워드 파일에서 텍스트 추출
def extract_text_from_word(file_content):
    doc = docx.Document(file_content)
    return '\n'.join([para.text for para in doc.paragraphs])

# PPT 파일에서 텍스트 추출
def extract_text_from_ppt(file_content):
    presentation = pptx.Presentation(file_content)
    text = ''
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + '\n'
    return text

# 이미지에서 텍스트 추출 (OCR)
def extract_text_from_image(file_content):
    image = Image.open(file_content)
    return "이미지에서 텍스트를 추출하는 기능은 구현되지 않았습니다."

# LLM을 통해 프롬프트와 파일을 전달하고 응답을 받는 함수
def run_llm_with_file_and_prompt(api_key, titles, requests, file_data_list):
    global global_generated_prompt
    openai.api_key = api_key

    responses = []
    global_generated_prompt = []  # 프롬프트들을 담을 리스트 초기화

    for i, (title, request, file_data) in enumerate(zip(titles, requests, file_data_list)):
        if isinstance(file_data, pd.DataFrame):
            file_data_str = file_data.to_string()
        elif isinstance(file_data, dict):
            file_data_str = "\n".join([f"시트 이름: {sheet_name}\n{data.to_string()}" for sheet_name, data in file_data.items()])
        else:
            file_data_str = str(file_data)

        generated_prompt = f"""
        보고서 제목은 '{title}'로 하고, 아래의 파일 데이터를 분석하여 '{request}'를 요구 사항을 만족할 수 있도록 최적화된 보고서를 완성해.
        파일 데이터: {file_data_str}
        """

        global_generated_prompt.append(generated_prompt)

        prompt_template = PromptTemplate(
            template=generated_prompt,
            input_variables=[]
        )

        llm = ChatOpenAI(model_name="gpt-4o")
        chain = LLMChain(llm=llm, prompt=prompt_template)
        
        success = False
        retry_count = 0
        max_retries = 5  # 최대 재시도 횟수
        
        while not success and retry_count < max_retries:
            try:
                response = chain.run({})
                responses.append(response)
                success = True  # 성공했을 때 True로 설정
            except RateLimitError:
                retry_count += 1
                st.warning(f"API 요청 한도를 초과했습니다. 10초 후 다시 시도합니다. 재시도 횟수: {retry_count}/{max_retries}")
                time.sleep(10)

            time.sleep(10)

    return responses

# 1 프레임: GitHub 정보 저장 및 OpenAI API 키 저장
st.subheader("1. GitHub 정보 저장 및 OpenAI API 키 저장")

col1, col2 = st.columns(2)

with col1:
    st.subheader("GitHub 정보 입력")
    repo_input = st.text_input("GitHub 저장소 (owner/repo 형식)", value="")
    branch_input = st.text_input("GitHub 브랜치", value="main")
    token_input = st.text_input("GitHub Token", type="password")

    if st.button("GitHub 정보 저장"):
        if repo_input and branch_input and token_input:
            st.session_state["github_repo"] = repo_input
            st.session_state["github_branch"] = branch_input
            st.session_state["github_token"] = token_input
            st.success("GitHub 정보가 성공적으로 저장되었습니다!")
        else:
            st.error("모든 GitHub 정보를 입력해야 합니다!")

with col2:
    st.subheader("OpenAI API 키 저장")
    api_key_input = st.text_input("OpenAI API 키를 입력하세요", type="password")
    if st.button("API 키 저장"):
        if api_key_input:
            st.session_state["api_key"] = api_key_input
            os.environ["OPENAI_API_KEY"] = api_key_input
            st.success("API 키가 성공적으로 저장되었습니다!")
        else:
            st.error("API 키를 입력해야 합니다!")

# 2 프레임: 파일 업로드
if not st.session_state.get("github_token") or not st.session_state.get("github_repo"):
    st.warning("GitHub 정보가 저장되기 전에는 파일 업로드를 할 수 없습니다. 먼저 GitHub 정보를 입력해 주세요.")
else:
    with st.expander("파일 업로드", expanded=True):
        uploaded_files = st.file_uploader("파일을 여러 개 드래그 앤 드롭하여 업로드하세요.", accept_multiple_files=True)

        if uploaded_files:
            for uploaded_file in uploaded_files:
                file_content = uploaded_file.read()
                file_name = uploaded_file.name
                folder_name = 'uploadFiles'

                sha = get_file_sha(st.session_state['github_repo'], f"{folder_name}/{file_name}", st.session_state['github_token'], branch=st.session_state['github_branch'])

                if sha:
                    st.warning(f"'{file_name}' 파일이 이미 존재합니다. 덮어쓰시겠습니까?")
                    col1, col2 = st.columns(2)

                    with col1:
                        if st.button(f"'{file_name}' 덮어쓰기", key=f"overwrite_{file_name}"):
                            upload_file_to_github(st.session_state['github_repo'], folder_name, file_name, file_content, st.session_state['github_token'], branch=st.session_state['github_branch'], sha=sha)
                            st.success(f"'{file_name}' 파일이 성공적으로 덮어쓰기 되었습니다.")
                            uploaded_files = None
                            break

                    with col2:
                        if st.button("취소", key=f"cancel_{file_name}"):
                            st.info("덮어쓰기가 취소되었습니다.")
                            uploaded_files = None
                            break
                else:
                    upload_file_to_github(st.session_state['github_repo'], folder_name, file_name, file_content, st.session_state['github_token'])
                    st.success(f"'{file_name}' 파일이 성공적으로 업로드되었습니다.")
                    uploaded_files = None

# 3 프레임: 작성 보고서 요청사항 테이블
st.subheader("3. 작성 보고서 요청사항 및 실행 버튼")

col1, col2 = st.columns([0.8, 0.2])

with col1:
    with st.expander("요청사항 리스트", expanded=True):
        if 'rows' not in st.session_state:
            st.session_state['rows'] = [{"제목": "", "요청": "", "파일": "", "데이터": "", "checked": False}]

        rows = st.session_state['rows']
        checked_rows = []

        for idx, row in enumerate(rows):
            with st.container():
                col1_1, col2_1 = st.columns([0.05, 0.95])
                with col1_1:
                    row_checked = st.checkbox(f"", key=f"row_checked_{idx}", value=row.get("checked", False), disabled=(idx == 0))
                with col2_1:
                    st.markdown(f"#### 요청사항 {idx+1}")

                    row['제목'] = st.text_input(f"제목 (요청사항 {idx+1})", row['제목'], key=f"title_{idx}")
                    row['요청'] = st.text_area(f"요청 (요청사항 {idx+1})", row['요청'], key=f"request_{idx}")

                    file_list = ['파일을 선택하세요.']
                    if st.session_state.get('github_token') and st.session_state.get('github_repo'):
                        file_list += get_github_files(st.session_state['github_repo'], st.session_state['github_branch'], st.session_state['github_token'], folder_name='uploadFiles')

                    selected_file = st.selectbox(f"파일 선택 (요청사항 {idx+1})", options=file_list, key=f"file_select_{idx}")

                    if selected_file != '파일을 선택하세요.':
                        file_path = selected_file
                        file_content = get_file_from_github(st.session_state["github_repo"], st.session_state["github_branch"], file_path, st.session_state["github_token"])
                        file_type = file_path.split('.')[-1].lower()

                        if file_type == 'xlsx':
                            file_data = extract_text_from_excel(file_content)
                        elif file_type == 'csv':
                            file_data = extract_text_from_csv(file_content)
                        elif file_type == 'docx':
                            file_data = extract_text_from_word(file_content)
                        elif file_type == 'pptx':
                            file_data = extract_text_from_ppt(file_content)
                        elif file_type == 'pdf':
                            file_data = extract_text_from_pdf(file_content)
                        elif file_type in ['png', 'jpg', 'jpeg']:
                            file_data = extract_text_from_image(file_content)
                        else:
                            st.error(f"지원되지 않는 파일 형식입니다: {file_type}")

                        row['파일'] = f"/{st.session_state['github_repo']}/{st.session_state['github_branch']}/uploadFiles/{selected_file}"
                        row['데이터'] = file_data

                    st.text_input(f"파일 경로 (요청사항 {idx+1})", row['파일'], disabled=True, key=f"file_{idx}")

                if row_checked:
                    checked_rows.append(idx)
                    row["checked"] = True
                else:
                    row["checked"] = False

        col1_1, col1_2, col1_3 = st.columns([0.33, 0.33, 0.33])

        with col1_1:
            if st.button("행 추가"):
                new_row = {"제목": "", "요청": "", "파일": "", "데이터": "", "checked": False}
                st.session_state['rows'].append(new_row)

        with col1_2:
            if st.button("행 삭제"):
                if checked_rows:
                    st.session_state['rows'] = [row for idx, row in enumerate(rows) if idx not in checked_rows]
                    st.success(f"체크된 {len(checked_rows)}개의 요청사항이 삭제되었습니다.")
                else:
                    st.warning("삭제할 요청사항을 선택해주세요.")

        with col1_3:
            if st.button("새로고침"):
                st.session_state['rows'] = st.session_state['rows']

with col2:
    st.write(" ")
    st.write(" ")
    if st.button("실행"):
        if not st.session_state.get("api_key"):
            st.error("먼저 OpenAI API 키를 입력하고 저장하세요!")
        elif not st.session_state['rows'] or all(not row["제목"] or not row["요청"] or not row["데이터"] for row in st.session_state['rows']):
            st.error("요청사항의 제목, 요청, 파일을 모두 입력해야 합니다!")
        else:
            titles = [row['제목'] for row in st.session_state['rows']]
            requests = [row['요청'] for row in st.session_state['rows']]
            file_data_list = [row['데이터'] for row in st.session_state['rows']]

            responses = run_llm_with_file_and_prompt(
                st.session_state["api_key"], 
                titles, 
                requests, 
                file_data_list
            )
            st.session_state["response"] = responses

# 4 프레임: 참고 탬플릿 보기 및 결과 보고서
st.subheader("4. 결과 보고서 및 5. 참고 탬플릿 보기")

col1, col2 = st.columns([0.5, 0.5])

with col1:
    st.subheader("5. 참고 탬플릿 보기")
    
    if not st.session_state.get("github_token") or not st.session_state.get("github_repo"):
        st.warning("GitHub 정보가 저장되기 전에는 템플릿 파일을 업로드하거나 선택할 수 없습니다. 먼저 GitHub 정보를 입력해 주세요.")
    else:
        uploaded_template = st.file_uploader("참고 템플릿 파일 업로드", type=["png", "jpg", "jpeg", "html", "pdf", "docx", "xlsx"])

        if uploaded_template:
            file_content = uploaded_template.read()
            file_name = uploaded_template.name
            folder_name = 'templateFiles'

            sha = get_file_sha(st.session_state['github_repo'], f"{folder_name}/{file_name}", st.session_state['github_token'], branch=st.session_state['github_branch'])

            if sha:
                st.warning(f"'{file_name}' 파일이 이미 존재합니다. 덮어쓰시겠습니까?")
                col1_1, col1_2 = st.columns([0.5, 0.5])

                with col1_1:
                    if st.button(f"'{file_name}' 덮어쓰기", key=f"overwrite_template_{file_name}"):
                        upload_file_to_github(st.session_state['github_repo'], folder_name, file_name, file_content, st.session_state['github_token'], branch=st.session_state['github_branch'], sha=sha)
                        st.success(f"'{file_name}' 파일이 성공적으로 덮어쓰기 되었습니다.")

                with col1_2:
                    if st.button("취소", key=f"cancel_template_{file_name}"):
                        st.info("덮어쓰기가 취소되었습니다.")
            else:
                upload_file_to_github(st.session_state['github_repo'], folder_name, file_name, file_content, st.session_state['github_token'])
                st.success(f"'{file_name}' 템플릿 파일이 성공적으로 업로드되었습니다.")

        template_files = get_github_files(st.session_state["github_repo"], st.session_state["github_branch"], st.session_state["github_token"], folder_name="templateFiles")
        
        selected_template = st.selectbox("탬플릿 리스트 선택", ["파일을 선택하세요."] + template_files)

        if st.button("저장"):
            if selected_template == "파일을 선택하세요.":
                st.error("파일을 선택하세요.")
            else:
                file_extension = selected_template.split('.')[-1].lower()
                template_url = f"https://raw.githubusercontent.com/{st.session_state['github_repo']}/blob/{st.session_state['github_branch']}/templateFiles/{selected_template}"
                template_relative_path = f"{st.session_state['github_repo']}/{st.session_state['github_branch']}/templateFiles/{selected_template}"
                st.session_state['selected_template_info'] = {
                    'url': template_url,
                    'relative_path': template_relative_path,
                    'extension': file_extension
                }

        if 'selected_template_info' in st.session_state:
            st.text_input("템플릿 파일 정보", st.session_state['selected_template_info']['relative_path'], disabled=True)

            # 템플릿 URL 정보 입력 추가
            st.text_input("템플릿 URL 정보", st.session_state['selected_template_info']['url'], disabled=True)

            file_extension = st.session_state['selected_template_info']['extension']
            file_url = st.session_state['selected_template_info']['url']

            if file_extension in ['png', 'jpg', 'jpeg', 'gif']:
                st.components.v1.html(f'<img src="{file_url}" alt="이미지 미리보기" style="max-width: 100%;">', height=400)
            elif file_extension == 'pdf':
                st.components.v1.html(f'<iframe src="{file_url}" width="100%" height="600px"></iframe>', height=600)
            elif file_extension == 'html':
                st.components.v1.html(f'<iframe src="{file_url}" width="100%" height="600px"></iframe>', height=600)
            elif file_extension == 'docx':
                st.write("Word 파일은 HTML 형식으로 변환하여 표시합니다.")
                file_content = get_file_from_github(st.session_state["github_repo"], st.session_state["github_branch"], f"templateFiles/{selected_template}", st.session_state["github_token"])
                word_data = extract_text_from_word(file_content)
                st.components.v1.html(f"<div>{word_data}</div>", height=600)
            elif file_extension == 'xlsx':
                st.write("Excel 파일은 HTML 형식으로 변환하여 표시합니다.")
                file_content = get_file_from_github(st.session_state["github_repo"], st.session_state["github_branch"], f"templateFiles/{selected_template}", st.session_state["github_token"])
                excel_data = extract_text_from_excel(file_content)
                excel_html = "".join([f"<h3>시트: {sheet_name}</h3>{data.to_html()}" for sheet_name, data in excel_data.items()])
                st.components.v1.html(f"<div>{excel_html}</div>", height=600)

with col2:
    st.subheader("4. 결과 보고서")
    st.text_area("전달된 프롬프트:", value="\n\n".join(global_generated_prompt), height=150)

    if "response" in st.session_state:
        for idx, response in enumerate(st.session_state["response"]):
            st.text_area(f"응답 {idx+1}:", value=response, height=300)
            st.components.v1.html(response, height=600, scrolling=True)
